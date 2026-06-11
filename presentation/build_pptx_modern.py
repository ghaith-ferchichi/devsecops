#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Soutenance — BTE Security AI Agent.
Thème « Modern Light Tech » (blanc + violet IA + teal + or, sans bold,
images/icônes, animations fluides). Réf : presentation/pptx-plan.md §3/§9/§11.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from PIL import Image

IMG = "/opt/devsecops/MASTER_PFE/img"
OUT = "/opt/devsecops/presentation/BTE_Security_AI_Agent_Soutenance.pptx"
GIFS = "/opt/devsecops/presentation/gifs"   # GIFs animés (auto-loop en mode Présentation)

# ---------- Palette : Modern Light Tech ----------
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
PANEL   = RGBColor(0xF2, 0xF4, 0xF7)
INK     = RGBColor(0x0B, 0x25, 0x45)
BODY    = RGBColor(0x33, 0x40, 0x5C)
VIOLET  = RGBColor(0x7C, 0x5C, 0xFC)   # AI fill
VIOLET_D= RGBColor(0x6D, 0x28, 0xD9)   # AI text
TEAL    = RGBColor(0x2B, 0xA6, 0xAE)   # tech fill
TEAL_D  = RGBColor(0x15, 0x7A, 0x82)   # tech text/lines
GOLD    = RGBColor(0xE0, 0xA1, 0x00)   # BTE fill
SLATE   = RGBColor(0x5B, 0x6B, 0x7B)
HAIR    = RGBColor(0xE1, 0xE7, 0xED)
GREEN   = RGBColor(0x1E, 0x9E, 0x5A)
AMBER   = RGBColor(0xE0, 0x86, 0x00)
RED     = RGBColor(0xD7, 0x26, 0x3D)
# light tints for decorative blobs
T_VIOLET= RGBColor(0xEF, 0xEA, 0xFE)
T_TEAL  = RGBColor(0xE6, 0xF6, 0xF7)
T_GOLD  = RGBColor(0xFBF1D8 >> 16 & 255, 0xFBF1D8 >> 8 & 255, 0xFBF1D8 & 255)

HEAD = "Segoe UI Semibold"
BODY_F = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
ANIM_BUILDS = []   # list of (slide, [shapes]) for fade-in entrance builds

# ===================== low-level helpers =====================
def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow: add_shadow(sp)
    return sp

def add_shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = spPr.makeelement(qn('a:outerShdw'),
        {'blurRad':'80000','dist':'30000','dir':'5400000','rotWithShape':'0'})
    clr = spPr.makeelement(qn('a:srgbClr'), {'val':'0B2545'})
    alpha = spPr.makeelement(qn('a:alpha'), {'val':'18000'})
    clr.append(alpha); sh.append(clr); el.append(sh); spPr.append(el)

def set_gradient(sp, c1, c2, angle_deg=0):
    spPr = sp._element.spPr
    for t in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        e = spPr.find(qn(t))
        if e is not None: spPr.remove(e)
    ang = int(angle_deg*60000)
    xml = ('<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
           '<a:gsLst>'
           f'<a:gs pos="0"><a:srgbClr val="{str(c1)}"/></a:gs>'
           f'<a:gs pos="100000"><a:srgbClr val="{str(c2)}"/></a:gs>'
           '</a:gsLst>'
           f'<a:lin ang="{ang}" scaled="1"/></a:gradFill>')
    grad = parse_xml(xml)
    ln = spPr.find(qn('a:ln'))
    if ln is not None: ln.addprevious(grad)
    else: spPr.append(grad)

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, sp_after=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        for (t, sz, col, bold, fnt, *rest) in para:
            ital = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = ital
            r.font.name = fnt; r.font.color.rgb = col
    return tb

def P(t, sz, col=BODY, bold=False, fnt=BODY_F, ital=False):
    return (t, sz, col, bold, fnt, ital)

def card_text(s, x, y, w, h, runs, fill=WHITE, line=HAIR, accent=None,
              shadow=True, anchor=MSO_ANCHOR.TOP, pad=0.28):
    """A single rounded-rect SHAPE whose own text frame holds `runs`
    (so it animates as ONE shape). Optional left accent bar drawn separately."""
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    if shadow: add_shadow(sp)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=Inches(pad); tf.margin_right=Inches(pad)
    tf.margin_top=Inches(0.18); tf.margin_bottom=Inches(0.14)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(3)
        for (t, sz, col, bold, fnt, *rest) in para:
            ital = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = ital
            r.font.name = fnt; r.font.color.rgb = col
    return sp

def fit_image(s, path, bx, by, bw, bh, card=True, caption=None):
    iw, ih = Image.open(path).size
    ar = iw/ih; bar = bw/bh
    if ar > bar: w = bw; h = bw/ar
    else: h = bh; w = bh*ar
    x = bx + (bw-w)/2; y = by + (bh-h)/2
    if card:
        pad = 0.10
        rect(s, x-pad, y-pad, w+2*pad, h+2*pad, fill=WHITE, line=HAIR,
             line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if caption:
        txt(s, bx, by+bh+0.05, bw, 0.3, [[P(caption, 10, SLATE, False, BODY_F, True)]],
            align=PP_ALIGN.CENTER)
    return (x, y, w, h)

def img(name):
    for sub in ("chapter_1","chapter_2","chapter_3","chapter_4","template"):
        p = os.path.join(IMG, sub, name)
        if os.path.exists(p): return p
    raise FileNotFoundError(name)

def gif_box(s, gif_name, bx, by, bw, bh, fallback=None, caption=None, label=None):
    """GIF animé (auto-loop en mode Présentation) cadré dans une carte.
    À défaut : image statique de repli + badge « ▶ GIF à venir » ; sinon placeholder."""
    gp = os.path.join(GIFS, gif_name)
    if os.path.exists(gp):
        return fit_image(s, gp, bx, by, bw, bh, card=True, caption=caption)
    if fallback:
        box = fit_image(s, img(fallback), bx, by, bw, bh, card=True, caption=caption)
        rect(s, bx+0.12, by+0.12, 1.45, 0.32, fill=VIOLET, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, bx+0.12, by+0.145, 1.45, 0.28, [[P("▶ GIF à venir", 9, WHITE, True, BODY_F)]], align=PP_ALIGN.CENTER)
        return box
    rect(s, bx, by, bw, bh, fill=PANEL, line=VIOLET, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, bx, by+bh/2-0.5, bw, 1.0, [[P("▶", 34, VIOLET, True, HEAD)],
        [P("GIF à venir — " + (label or gif_name), 12, SLATE, False, BODY_F)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return (bx, by, bw, bh)

# ===================== decorative + icon helpers =====================
def chevron(s, x, y, size, color, n=3, gap=0.18):
    for i in range(n):
        c = rect(s, x+i*gap, y, size, size, fill=color, shape=MSO_SHAPE.CHEVRON)
    return

def dotted_grid(s, x, y, cols, rows, gap=0.22, d=0.05, color=HAIR):
    for r in range(rows):
        for c in range(cols):
            rect(s, x+c*gap, y+r*gap, d, d, fill=color, shape=MSO_SHAPE.OVAL)

def blob(s, x, y, w, h, color, shape=MSO_SHAPE.OVAL):
    rect(s, x, y, w, h, fill=color, shape=shape)

def arcs(s, cx, cy, color=HAIR, n=3, r0=0.5, step=0.35):
    for i in range(n):
        r = r0 + i*step
        rect(s, cx-r, cy-r, 2*r, 2*r, fill=None, line=color, line_w=1.2, shape=MSO_SHAPE.OVAL)

def deco(s, kind="light"):
    """Subtle background decoration. Call right after the white canvas."""
    if kind == "light":
        blob(s, 11.9, -0.7, 1.9, 1.9, T_VIOLET)
        dotted_grid(s, 11.0, 0.35, 6, 3, color=HAIR)
        blob(s, -0.6, 6.0, 1.7, 1.7, T_TEAL)
    elif kind == "divider":
        blob(s, 11.4, -0.9, 2.6, 2.6, T_VIOLET)
        blob(s, -0.9, 5.6, 2.4, 2.4, T_TEAL)
        arcs(s, 12.7, 7.4, color=RGBColor(0xEDE9FB>>16&255,0xEDE9FB>>8&255,0xEDE9FB&255), n=3, r0=0.5)
    elif kind == "hero":
        dotted_grid(s, 0.7, 0.35, 6, 3, color=HAIR)
        blob(s, 11.9, 5.9, 1.9, 1.9, T_GOLD)

def icon_badge(s, x, y, d, label, color):
    """Circular badge with a short text glyph/number (geometric, swap later)."""
    rect(s, x, y, d, d, fill=color, shape=MSO_SHAPE.OVAL)
    txt(s, x, y+d*0.16, d, d*0.7, [[P(label, 18, WHITE, True, HEAD)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===================== animation injection =====================
def add_fade_transition(sld_obj, spd="med"):
    sld = sld_obj._element
    for t in sld.findall(qn('p:transition')): sld.remove(t)
    xml = ('<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
           f'spd="{spd}"><p:fade/></p:transition>')
    el = parse_xml(xml)
    anchor = sld.find(qn('p:clrMapOvr'))
    if anchor is None: anchor = sld.find(qn('p:cSld'))
    sld.insert(list(sld).index(anchor)+1, el)

def add_fade_in_build(sld_obj, shapes, dur=500):
    """Standard mainSeq click-to-reveal fade-in; one click per shape."""
    sld = sld_obj._element
    for t in sld.findall(qn('p:timing')): sld.remove(t)
    ids = [sp.shape_id for sp in shapes]
    cid = [2]
    def nid():
        cid[0]+=1; return cid[0]
    steps = []
    for spid in ids:
        a,b,c,d,e = nid(),nid(),nid(),nid(),nid()
        steps.append(
          f'<p:par><p:cTn id="{a}" fill="hold"><p:stCondLst><p:cond delay="indefinite"/>'
          f'</p:stCondLst><p:childTnLst>'
          f'<p:par><p:cTn id="{b}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
          f'<p:childTnLst>'
          f'<p:par><p:cTn id="{c}" presetID="10" presetClass="entr" presetSubtype="0" '
          f'fill="hold" nodeType="clickEffect"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
          f'<p:childTnLst>'
          f'<p:set><p:cBhvr><p:cTn id="{d}" dur="1" fill="hold"><p:stCondLst>'
          f'<p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
          f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>'
          f'<p:to><p:strVal val="visible"/></p:to></p:set>'
          f'<p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{e}" dur="{dur}"/>'
          f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>'
          f'</p:childTnLst></p:cTn></p:par>'
          f'</p:childTnLst></p:cTn></p:par>'
          f'</p:childTnLst></p:cTn></p:par>')
    blds = "".join(f'<p:bldP spid="{i}" grpId="0" animBg="1"/>' for i in ids)
    xml = ('<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
           '<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
           '<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
           '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
           + "".join(steps) +
           '</p:childTnLst></p:cTn>'
           '<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
           '<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
           '</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
           f'<p:bldLst>{blds}</p:bldLst></p:timing>')
    el = parse_xml(xml)
    anchor = sld.find(qn('p:transition'))
    if anchor is None: anchor = sld.find(qn('p:clrMapOvr'))
    if anchor is None: anchor = sld.find(qn('p:cSld'))
    sld.insert(list(sld).index(anchor)+1, el)

# ===================== composite components =====================
_PAGE = [0]   # numérotation automatique des slides de contenu
def header(s, kicker, title, num=True, deco_kind="light"):
    rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
    deco(s, deco_kind)
    if kicker:
        txt(s, 0.7, 0.5, 11.0, 0.35, [[P(kicker.upper(), 12, TEAL_D, True, BODY_F)]])
    txt(s, 0.7, 0.8, 12.0, 0.9, [[P(title, 28, INK, False, HEAD)]])
    rect(s, 0.72, 1.62, 0.5, 0.07, fill=VIOLET, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 1.30, 1.63, 1.0, 0.05, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if num:
        _PAGE[0] += 1
        footer(s, _PAGE[0])
    else:
        footer(s, None)

def footer(s, num=None):
    rect(s, 0.7, 7.04, 11.93, 0.012, fill=HAIR)
    txt(s, 0.7, 7.1, 8.0, 0.3,
        [[P("BTE Security AI Agent", 9, SLATE, False, BODY_F),
          P("  ·  Ghaith Ferchichi", 9, SLATE, False, BODY_F)]])
    if num is not None:
        chevron(s, 11.55, 7.12, 0.12, VIOLET, n=2, gap=0.11)
        txt(s, 11.9, 7.07, 0.7, 0.3, [[P(f"{num:02d}", 10, INK, True, HEAD)]], align=PP_ALIGN.RIGHT)

PLAN_NODES = ["Contexte", "État de l'art", "Conception", "Réalisation", "Conclusion"]
def plan_divider(s, active, y=3.5):
    n = len(PLAN_NODES); x0 = 1.2; x1 = 12.13; d = 0.62
    span = x1 - x0; gap = span/(n-1)
    rect(s, x0+d/2, y+d/2-0.01, span, 0.02, fill=HAIR)
    for i, lab in enumerate(PLAN_NODES):
        cx = x0 + i*gap
        on = (i == active)
        c = VIOLET if on else PANEL
        dd = d if on else d*0.82
        rect(s, cx+(d-dd)/2, y+(d-dd)/2, dd, dd, fill=c,
             line=(None if on else HAIR), line_w=1.2, shape=MSO_SHAPE.OVAL)
        txt(s, cx+(d-dd)/2, y+(d-dd)/2+dd*0.16, dd, dd*0.7,
            [[P(str(i+1), 16 if on else 13, WHITE if on else SLATE, True, HEAD)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, cx-0.6, y+d+0.12, d+1.2, 0.4,
            [[P(lab, 12 if on else 10.5, INK if on else SLATE, on, BODY_F)]],
            align=PP_ALIGN.CENTER)

def section_divider(s, active, subtitle):
    rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
    deco(s, "divider")
    num = active+1
    txt(s, 0.7, 0.9, 4.0, 1.6, [[P(f"{num:02d}", 78, RGBColor(0xE7,0xE3,0xF7), True, HEAD)]])
    txt(s, 0.72, 2.0, 11.0, 0.9, [[P(PLAN_NODES[active], 40, INK, False, HEAD)]])
    rect(s, 0.74, 2.92, 0.5, 0.08, fill=VIOLET, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 1.32, 2.93, 1.2, 0.06, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 0.74, 3.12, 11.0, 0.5, [[P(subtitle, 15, SLATE, False, BODY_F, True)]])
    plan_divider(s, active, y=4.7)
    footer(s, None)

# ===================== SLIDES =====================

# ---- 1. PAGE DE TITRE ----
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
# bandeau hero (dégradé) + image, à droite
panel = rect(s, 8.25, 0, 5.083, 7.5, fill=VIOLET); set_gradient(panel, VIOLET, TEAL, 60)
fit_image(s, img("solution_overview.png"), 8.5, 2.6, 4.6, 2.6, card=True)
chevron(s, 8.6, 5.85, 0.45, WHITE, n=3, gap=0.4)
# logos institutionnels : BTE · ISI · UTM
def _logo(nm, x, y, h):
    try: s.shapes.add_picture(img(nm), Inches(x), Inches(y), height=Inches(h))
    except Exception: pass
_logo("logo_BTE.png", 0.7, 0.5, 0.55)
_logo("LogoISI.png", 5.0, 0.5, 0.55)
_logo("Logo_UTM.png", 6.95, 0.46, 0.62)
txt(s, 0.7, 1.85, 7.3, 0.35, [[P("UNIVERSITÉ DE TUNIS EL MANAR · INSTITUT SUPÉRIEUR D'INFORMATIQUE", 9.5, SLATE, True, BODY_F)]])
txt(s, 0.7, 2.2, 7.2, 0.35, [[P("Rapport de Stage de Fin d'Études", 13, TEAL_D, True, HEAD)]])
txt(s, 0.7, 2.62, 7.3, 1.5, [
    [P("BTE Security", 36, INK, False, HEAD)],
    [P("AI Agent", 36, VIOLET_D, False, HEAD)],
])
txt(s, 0.7, 3.98, 7.3, 0.6, [[P("Agent IA pour la revue automatisée de code",
                                14, BODY, False, BODY_F)]])
txt(s, 0.7, 4.55, 7.4, 0.7, [
    [P("Diplôme National de Mastère Professionnel", 12.5, INK, True, HEAD)],
    [P("Spécialité : Sécurité des Systèmes d'Information et des Infrastructures", 11.5, BODY, False, BODY_F)],
])
rect(s, 0.72, 5.42, 1.3, 0.05, fill=VIOLET)
txt(s, 0.7, 5.58, 7.4, 1.5, [
    [P("Réalisé par  ", 12, SLATE, False, BODY_F), P("Ghaith FERCHICHI", 12, INK, True, HEAD)],
    [P("Encadrant professionnel : ", 10.5, SLATE, False, BODY_F), P("M. Kamel KAOUECH", 10.5, BODY, False, BODY_F)],
    [P("Encadrante académique : ", 10.5, SLATE, False, BODY_F), P("Mme Ghayet El Mouna ZHIOUA", 10.5, BODY, False, BODY_F)],
    [P("Banque de Tunisie et des Émirats · Année universitaire 2025–2026", 10.5, SLATE, False, BODY_F, True)],
])

# ---- 2. SOMMAIRE ----
s = slide(); header(s, "Plan", "Sommaire", None)
somm = [("1","Contexte général","Organisme d'accueil, étude de l'existant, problématique et méthodologie"),
        ("2","État de l'art et choix technologiques","Concepts DevSecOps, spécification des besoins et technologies retenues"),
        ("3","Conception","Architecture de l'agent — « un cerveau, et non un simple pipeline »"),
        ("4","Réalisation et résultats","Sécurité, démonstration, observabilité et validation"),
        ("5","Conclusion et perspectives","Limites, perspectives et bilan du projet")]
yy=1.95
for i,(n,t,d) in enumerate(somm):
    c = VIOLET if i % 2 == 0 else TEAL
    icon_badge(s, 0.9, yy, 0.62, n, c)
    txt(s, 1.75, yy+0.04, 10.7, 0.7, [[P(t, 16, INK, True, HEAD)],[P(d, 12, SLATE, False, BODY_F)]])
    yy += 0.98

# ---- 3. INTRODUCTION ----
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE); deco(s, "light")
blob(s, 9.6, 4.6, 3.4, 3.4, T_VIOLET)
chevron(s, 0.9, 1.4, 0.34, VIOLET, n=3, gap=0.3)
txt(s, 0.9, 2.0, 11.3, 0.5, [[P("INTRODUCTION", 14, TEAL_D, True, BODY_F)]])
txt(s, 0.9, 2.5, 11.5, 2.0, [
    [P("Une seule injection SQL fusionnée dans une API de paiement", 30, INK, False, HEAD)],
    [P("peut compromettre l'intégrité du système d'information de la banque.", 28, VIOLET_D, False, HEAD)],
])
rect(s, 0.92, 4.95, 1.6, 0.05, fill=GOLD)
txt(s, 0.9, 5.2, 11.0, 1.0, [
    [P("À la BTE, la revue de sécurité demeure ", 17, BODY, False, BODY_F),
     P("manuelle", 17, GOLD, True, HEAD),
     P(", peut atteindre ", 17, BODY, False, BODY_F),
     P("vingt-quatre heures", 17, GOLD, True, HEAD),
     P(",", 17, BODY, False, BODY_F)],
    [P("et dépend de la disponibilité et de l'expertise du relecteur.", 17, BODY, False, BODY_F)],
])

# ======================= CHAPITRE 1 — CONTEXTE GÉNÉRAL =======================
s = slide(); section_divider(s, 0, "Organisme d'accueil, étude de l'existant et méthodologie de travail")

# ---- Organisme & étude de l'existant ----
s = slide(); header(s, "01 · Contexte général", "Organisme d'accueil et étude de l'existant", 1)
fit_image(s, img("asis_workflow.png"), 0.7, 1.85, 6.0, 5.0, caption="Processus de revue actuel (AS-IS) à la BTE")
txt(s, 7.1, 2.0, 5.5, 0.5, [[P("Banque de Tunisie et des Émirats", 16, INK, True, HEAD)]])
for i,(t,d) in enumerate([("Banque universelle", "particuliers, professionnels, entreprises"),
                          ("DCIO — sécurité opérationnelle", "cadre d'accueil du stage"),
                          ("Convention État tunisien · ADIA", "création en 1982, vocation élargie")]):
    rect(s, 7.1, 2.7+i*0.78, 0.16, 0.16, fill=TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, 7.4, 2.62+i*0.78, 5.2, 0.7, [[P(t+" — ", 13, INK, True, HEAD), P(d, 12, BODY, False, BODY_F)]])
card_text(s, 7.1, 5.2, 5.5, 1.3, [
    [P("Une revue intégralement manuelle", 14, RED, True, HEAD)],
    [P("aucune analyse statique ni gate CI/CD n'est greffé sur le cycle de livraison.", 12, BODY, False, BODY_F)],
], fill=PANEL, line=None)

# ---- Problématique ----
s = slide(); header(s, "01 · Contexte général", "Problématique", 2)
rect(s, 0.7, 2.0, 11.93, 4.4, fill=PANEL, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
items = [("Délai pouvant atteindre 24 h", "selon la disponibilité et l'expertise du relecteur", RED),
         ("Absence de gate CI/CD", "une Pull Request vulnérable peut être fusionnée sans blocage", RED),
         ("Qualité de revue variable", "la couverture de sécurité diffère d'un relecteur à l'autre", AMBER),
         ("Angles morts", "secrets commités et dépendances vulnérables non détectés à l'œil nu", AMBER)]
for i,(t,d,c) in enumerate(items):
    col = i % 2; row = i // 2
    x = 1.1 + col*5.9; y = 2.5 + row*1.9
    icon_badge(s, x, y, 0.7, "!", c)
    txt(s, x+0.95, y+0.02, 4.7, 1.4, [[P(t, 15, INK, True, HEAD)],[P(d, 12.5, BODY, False, BODY_F)]])

# ---- Objectifs et contribution (stats animées) ----
s = slide(); header(s, "01 · Contexte général", "Objectifs et contribution", 3)
card_text(s, 0.7, 1.95, 11.93, 1.75, [
    [P("Un ", 19, BODY, False, BODY_F), P("agent IA autonome", 19, VIOLET_D, True, HEAD),
     P(" qui revoit chaque Pull Request selon l'", 19, BODY, False, BODY_F),
     P("OWASP Top 10", 19, INK, True, HEAD), P(" en ", 19, BODY, False, BODY_F),
     P("quinze minutes", 19, INK, True, HEAD), P(",", 19, BODY, False, BODY_F)],
    [P("publie les correctifs sur GitHub, et s'exécute ", 19, BODY, False, BODY_F),
     P("intégralement en local", 19, GREEN, True, HEAD),
     P(" — aucun code ne quitte la banque.", 19, BODY, False, BODY_F)],
], fill=PANEL, line=None, anchor=MSO_ANCHOR.MIDDLE)
rect(s, 0.7, 1.95, 0.13, 1.75, fill=VIOLET, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
stats = [("24 h → 15 min", "délai de revue", VIOLET_D),
         ("5 / 5", "vulnérabilités détectées", GREEN),
         ("0", "ligne de code ne sort du VPS", TEAL_D)]
cards = []
for i,(big_t,lab,c) in enumerate(stats):
    x = 0.7 + i*4.06
    sp = card_text(s, x, 4.15, 3.8, 2.0, [
        [P(big_t, 30, c, True, HEAD)], [P("", 6, BODY)], [P(lab, 13, SLATE, False, BODY_F)],
    ], anchor=MSO_ANCHOR.MIDDLE)
    for pp in sp.text_frame.paragraphs: pp.alignment = PP_ALIGN.CENTER
    rect(s, x, 4.15, 3.8, 0.12, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    cards.append(sp)
ANIM_BUILDS.append((s, cards))

# ---- Solution proposée (TO-BE) ----
s = slide(); header(s, "01 · Contexte général", "Solution proposée (TO-BE)", 4)
fit_image(s, img("tobe_workflow.png"), 0.7, 1.8, 6.6, 5.05, caption="Processus cible (TO-BE) — BTE Security AI Agent")
txt(s, 7.7, 2.1, 5.0, 0.6, [[P("Un agent qui prend en charge tout le cycle", 16, INK, True, HEAD)]])
yy = 2.8
for t in ["Webhook signé, pipeline automatique","Cinq scanners SAST et deux LLM locaux",
          "Revue OWASP, score et verdict","Gate CI/CD (APPROVE / CHANGES / BLOCK)","Inférence locale — conformité BCT"]:
    chevron(s, 7.7, yy+0.03, 0.16, VIOLET, n=2, gap=0.13)
    txt(s, 8.1, yy-0.04, 4.6, 0.5, [[P(t, 13, BODY, False, BODY_F)]]); yy += 0.62
card_text(s, 7.7, 6.0, 5.0, 0.7, [[P("Délai : 15–25 min  vs  jusqu'à 24 h en manuel", 13, GREEN, True, HEAD)]],
          fill=PANEL, line=None, anchor=MSO_ANCHOR.MIDDLE)

# ---- Méthodologie ----
s = slide(); header(s, "01 · Contexte général", "Méthodologie de travail", 5)
fit_image(s, img("gantt_chart.png"), 0.7, 1.95, 7.7, 4.4, caption="Diagramme de Gantt — dix sprints sur cinq mois")
txt(s, 8.7, 2.05, 4.0, 0.5, [[P("Méthode Scrumban", 16, INK, True, HEAD)]])
for t,d in [("Scrum + Kanban","cadence de sprints et flux à WIP = 1"),
            ("Dix sprints / 5 mois","de la fondation à l'extension"),
            ("Système pull","chaque incident devient la carte prioritaire suivante")]:
    pass
yy=2.7
for t,d in [("Scrum + Kanban","cadence de sprints, flux à WIP = 1"),
            ("Dix sprints sur cinq mois","de la fondation Docker à l'observabilité"),
            ("Système « pull »","chaque incident de production devient la carte prioritaire")]:
    rect(s, 8.7, yy+0.04, 0.16,0.16, fill=VIOLET, shape=MSO_SHAPE.OVAL)
    txt(s, 9.0, yy-0.04, 3.8, 0.8, [[P(t, 13, INK, True, HEAD)],[P(d, 11.5, BODY, False, BODY_F)]]); yy += 1.0

# ======================= CHAPITRE 2 — ÉTAT DE L'ART =======================
s = slide(); section_divider(s, 1, "Concepts DevSecOps, spécification des besoins et choix technologiques")

# ---- DevSecOps & shift-left ----
s = slide(); header(s, "02 · État de l'art", "DevSecOps et approche shift-left", 6)
fit_image(s, img("devsecops_lifecycle.jpg"), 0.7, 1.9, 6.4, 4.6, caption="Cycle de vie DevSecOps")
txt(s, 7.4, 2.1, 5.2, 0.5, [[P("Intégrer la sécurité au plus tôt", 16, INK, True, HEAD)]])
for t,d in [("Security as Code","politiques de sécurité versionnées et automatisées"),
            ("Shift Left","détecter la vulnérabilité dès la Pull Request"),
            ("Le verdict devient un gate","contrôle technique conditionnant la fusion")]:
    pass
yy=2.75
for t,d in [("Security as Code","politiques de sécurité versionnées et automatisées"),
            ("Shift Left","détecter la vulnérabilité dès la Pull Request"),
            ("Le verdict comme gate","un contrôle technique conditionne la fusion")]:
    rect(s, 7.4, yy+0.04, 0.16,0.16, fill=TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, 7.7, yy-0.04, 4.9, 0.85, [[P(t, 13.5, INK, True, HEAD)],[P(d, 12, BODY, False, BODY_F)]]); yy += 1.05

# ---- Spécification des besoins ----
s = slide(); header(s, "02 · État de l'art", "Spécification des besoins", 7)
fp = card_text(s, 0.7, 2.0, 5.85, 4.4, [
    [P("Besoins fonctionnels", 16, VIOLET_D, True, HEAD)], [P("", 4, BODY)],
    [P("•  Détecter les vulnérabilités (5 scanners)", 13, BODY, False, BODY_F)],
    [P("•  Produire une revue OWASP + score + verdict", 13, BODY, False, BODY_F)],
    [P("•  Publier les commentaires inline sur GitHub", 13, BODY, False, BODY_F)],
    [P("•  Conditionner la fusion (gate CI/CD)", 13, BODY, False, BODY_F)],
    [P("•  Assistant conversationnel d'exploitation", 13, BODY, False, BODY_F)],
], fill=PANEL, line=None)
fp = card_text(s, 6.78, 2.0, 5.85, 4.4, [
    [P("Besoins non fonctionnels", 16, TEAL_D, True, HEAD)], [P("", 4, BODY)],
    [P("•  Confidentialité — 100 % local (conformité BCT)", 13, BODY, False, BODY_F)],
    [P("•  Performance — 15 à 25 min par revue", 13, BODY, False, BODY_F)],
    [P("•  Résilience — reprise et circuit breaker", 13, BODY, False, BODY_F)],
    [P("•  Observabilité — métriques et alertes", 13, BODY, False, BODY_F)],
    [P("•  Sécurité de la plateforme — HMAC, secrets exclus", 13, BODY, False, BODY_F)],
], fill=T_TEAL, line=None)

# ---- Choix des modèles LLM — benchmark ----
s = slide(); header(s, "02 · État de l'art", "Choix des modèles LLM — benchmark sur le matériel cible", 8)
rows = [("Modèle","Taille","Vitesse","Précision outils","Décision"),
        ("qwen2.5-coder:7b","4,7 Go","~5 tok/s","80 %","Classification + chat"),
        ("qwen2.5-coder:14b","9,0 Go","~3,2 tok/s","80 %","Analyse de sécurité"),
        ("llama3.2:3b","2,0 Go","~8 tok/s","0 %","Écarté"),
        ("granite3.1-dense:2b","1,6 Go","~8,5 tok/s","0 %","Écarté")]
tb = s.shapes.add_table(len(rows),5,Inches(0.7),Inches(1.95),Inches(11.93),Inches(3.1)).table
tb.columns[0].width=Inches(3.3); tb.columns[1].width=Inches(1.7); tb.columns[2].width=Inches(1.9)
tb.columns[3].width=Inches(2.4); tb.columns[4].width=Inches(2.63)
for r in range(len(rows)):
    for c in range(5):
        cell = tb.cell(r,c); cell.margin_left=Inches(0.12); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
        run = cell.text_frame.paragraphs[0].add_run(); run.text = rows[r][c]
        run.font.name = HEAD if (r==0 or c==0) else BODY_F
        run.font.size = Pt(12.5)
        if r==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=INK; run.font.color.rgb=WHITE; run.font.bold=True
        elif r in (1,2):
            cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0xEA,0xF7,0xF0)
            run.font.color.rgb = GREEN if c==4 else INK; run.font.bold=(c==0 or c==4)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r%2 else PANEL
            run.font.color.rgb = BODY
            if c==4: run.font.color.rgb = RED
card_text(s, 0.7, 5.35, 7.6, 1.25, [
    [P("Architecture à deux modèles : ", 13, INK, True, HEAD),
     P("le 7B classifie en ~30 s, le 14B mène l'analyse approfondie — "
       "à précision d'appel d'outils égale (80 %).", 13, BODY, False, BODY_F)],
], fill=PANEL, line=None, anchor=MSO_ANCHOR.MIDDLE)
card_text(s, 8.5, 5.35, 4.13, 1.25, [
    [P("Ollama : +22 %", 16, TEAL_D, True, HEAD)],
    [P("de débit vs LocalAI, sur GGUF identique (benchmark croisé).", 11.5, BODY, False, BODY_F)],
], fill=T_TEAL, line=None, anchor=MSO_ANCHOR.MIDDLE)

# ---- Technologies retenues (hub) ----
s = slide(); header(s, "02 · État de l'art", "Technologies retenues", 9)
triad = [("5", "Cinq scanners SAST", "Trivy · Gitleaks · Semgrep · Checkov · OSV", TEAL),
         ("IA", "Deux LLM locaux", "qwen2.5-coder 7B (classer) + 14B (analyser)", VIOLET),
         ("⟳", "Orchestration", "LangGraph — graphe d'état, reprise, routage", GOLD)]
for i,(g,t,d,c) in enumerate(triad):
    x = 0.95 + i*4.05
    rect(s, x, 2.1, 3.7, 3.7, fill=WHITE, line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    icon_badge(s, x+1.45, 2.5, 0.8, g, c)
    txt(s, x+0.3, 3.6, 3.1, 0.5, [[P(t, 16, INK, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, x+0.3, 4.2, 3.1, 1.3, [[P(d, 12.5, BODY, False, BODY_F)]], align=PP_ALIGN.CENTER)
txt(s, 0.7, 6.1, 11.93, 0.5, [[P("Sortie SAST nettoyée (−52 % de tokens) avant transmission au LLM.",
                                 13, SLATE, False, BODY_F, True)]], align=PP_ALIGN.CENTER)

# ---- Positionnement ----
s = slide(); header(s, "02 · État de l'art", "Positionnement par rapport aux solutions existantes", 9)
rows = [("Solution","Données / inférence","LLM","100 % local"),
        ("GitHub CodeQL","Cloud GitHub","Non","Non"),
        ("Snyk","Code envoyé au cloud","Limité","Non"),
        ("SonarQube","Auto-hébergeable","Non","Règles statiques"),
        ("Copilot Autofix","Cloud","Oui","Non"),
        ("BTE Security AI Agent","Rien ne sort du VPS","Oui","OUI")]
tb = s.shapes.add_table(len(rows),4,Inches(0.7),Inches(1.95),Inches(11.93),Inches(3.7)).table
tb.columns[0].width=Inches(3.5); tb.columns[1].width=Inches(4.4)
tb.columns[2].width=Inches(1.6); tb.columns[3].width=Inches(2.43)
for r in range(len(rows)):
    for c in range(4):
        cell = tb.cell(r,c); cell.margin_left=Inches(0.12); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
        run = cell.text_frame.paragraphs[0].add_run(); run.text = rows[r][c]
        run.font.name = HEAD if (r==0 or c==0) else BODY_F
        run.font.size = Pt(12.5 if r else 13)
        if r==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=INK; run.font.color.rgb=WHITE; run.font.bold=True
        elif r==len(rows)-1:
            cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0xEA,0xF7,0xF0)
            run.font.color.rgb = GREEN if c==3 else INK; run.font.bold=(c==0 or c==3)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r%2 else PANEL
            run.font.color.rgb = BODY
            if c==3 and rows[r][c]=="Non": run.font.color.rgb = RED
txt(s, 0.7, 5.85, 11.93, 0.5, [[P("Différenciateur : ", 14, INK, True, HEAD),
    P("l'inférence 100 % locale garantit la confidentialité bancaire et la conformité BCT.", 14, GREEN, True, HEAD)]])

# ======================= CHAPITRE 3 — CONCEPTION =======================
s = slide(); section_divider(s, 2, "Architecture de l'agent — « un cerveau, et non un simple pipeline »")

# ---- Architecture globale ----
s = slide(); header(s, "03 · Conception", "Architecture globale de la solution", 10)
fit_image(s, img("full_architecture.png"), 0.7, 1.8, 8.5, 5.0, caption="Onze conteneurs répartis en quatre couches")
layers = [("Couche entrée","nginx — webhook HMAC",TEAL),("Couche IA","FastAPI + LangGraph + Ollama",VIOLET),
          ("Couche données","PostgreSQL + Redis",INK),("Observabilité","Prometheus · Grafana",GREEN)]
for i,(t,d,c) in enumerate(layers):
    y = 2.2 + i*1.1
    rect(s, 9.5, y, 0.13, 0.85, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 9.8, y, 3.1, 0.9, [[P(t, 14, INK, True, HEAD)],[P(d, 11, BODY, False, BODY_F)]])

# ---- Manifeste 6 capacités (animées) ----
s = slide(); header(s, "03 · Conception", "Un agent autonome, et non un simple appel LLM", 11)
caps = [("1","Perception","événements webhook GitHub",TEAL),
        ("2","Raisonnement","classe la PR, route le graphe d'état",VIOLET),
        ("3","Action","commente, fixe le gate, escalade",INK),
        ("4","État et reprise","checkpoints LangGraph (PostgreSQL)",GREEN),
        ("5","Usage d'outils","dix-neuf outils réels (Docker, Prometheus…)",VIOLET),
        ("6","Auto-exploitation","scheduler, health digest, gardien disque",TEAL)]
cards=[]
for i,(n,t,d,c) in enumerate(caps):
    col=i%3; row=i//3
    x=0.7+col*4.06; y=1.95+row*2.35
    sp = card_text(s, x, y, 3.8, 2.05, [
        [P(n+"  ", 18, c, True, HEAD), P(t, 15, INK, True, HEAD)], [P("", 4, BODY)],
        [P(d, 12.5, BODY, False, BODY_F)],
    ])
    cards.append(sp)
ANIM_BUILDS.append((s, cards))

# ---- Pipeline (image + 7 étapes animées) ----
s = slide(); header(s, "03 · Conception", "Le pipeline de revue — graphe d'état LangGraph", 12)
gif_box(s, "pipeline.gif", 0.7, 1.8, 7.0, 5.0, fallback="langgraph_state_graph.png", caption="StateGraph à neuf nœuds, routage conditionnel", label="exécution du pipeline")
steps = ["webhook — HMAC + déduplication","classify — LLM 7B (~30 s)","route — scanners selon la PR",
         "scan — cinq SAST en parallèle","analyze — LLM 14B + OWASP","verdict — APPROVE/CHANGES/BLOCK",
         "publish — commentaires + gate"]
step_shapes=[]
yy=2.0
for i,st in enumerate(steps):
    sp = card_text(s, 8.0, yy, 4.7, 0.62, [
        [P(f"{i+1}  ", 13, (VIOLET_D if i in (1,4) else TEAL_D), True, HEAD),
         P(st, 12, BODY, False, BODY_F)]], anchor=MSO_ANCHOR.MIDDLE, pad=0.18)
    step_shapes.append(sp); yy += 0.68
ANIM_BUILDS.append((s, step_shapes))

# ---- Fiabilité : anti-hallucination + 2 modèles ----
s = slide(); header(s, "03 · Conception", "Fiabilité — anti-hallucination et deux modèles", 13)
fit_image(s, img("anti_hallucination_layers.png"), 0.7, 1.8, 6.0, 4.5, caption="Six couches de protection")
fit_image(s, img("two_model_architecture.png"), 7.0, 1.8, 5.6, 3.0, caption="Deux modèles : 7B pour classer, 14B pour analyser")
card_text(s, 7.0, 5.1, 5.6, 1.5, [
    [P("Répondre au doute n°1 du jury sur les LLM", 13, INK, True, HEAD)],
    [P("temperature=0 · num_ctx maîtrisé · garde sans-outil · prompt anti-hallucination · "
       "parser éliminant les lignes inexistantes · fusion réduisant le temps de moitié.", 12, BODY, False, BODY_F)],
], fill=PANEL, line=None)

# ======================= CHAPITRE 4 — RÉALISATION =======================
s = slide(); section_divider(s, 3, "Environnement, sécurité en action, démonstration et validation")

# ---- Environnement de travail ----
s = slide(); header(s, "04 · Réalisation", "Environnement de travail")
card_text(s, 0.7, 2.0, 5.85, 4.4, [
    [P("Environnement matériel — VPS de production", 15, INK, True, HEAD)], [P("", 5, BODY)],
    [P("•  CPU : 12 cœurs Intel Haswell (AVX2)", 13, BODY, False, BODY_F)], [P("", 3, BODY)],
    [P("•  RAM : 45 Go  ·  Disque : 290 Go", 13, BODY, False, BODY_F)], [P("", 3, BODY)],
    [P("•  GPU : aucun — inférence CPU uniquement", 13, RED, True, BODY_F)], [P("", 3, BODY)],
    [P("•  OS : Ubuntu Linux", 13, BODY, False, BODY_F)],
], fill=PANEL, line=None)
card_text(s, 6.78, 2.0, 5.85, 4.4, [
    [P("Environnement logiciel — pile open source", 15, INK, True, HEAD)], [P("", 5, BODY)],
    [P("•  Docker 29.4 — onze conteneurs", 13, BODY, False, BODY_F)], [P("", 3, BODY)],
    [P("•  Python 3.12 · FastAPI · LangGraph 1.1+", 13, BODY, False, BODY_F)], [P("", 3, BODY)],
    [P("•  PostgreSQL 16 · Redis 7 · nginx", 13, BODY, False, BODY_F)], [P("", 3, BODY)],
    [P("•  Trivy · Gitleaks · Semgrep · Checkov · OSV-Scanner", 13, BODY, False, BODY_F)],
], fill=T_TEAL, line=None)
txt(s, 0.7, 6.55, 11.93, 0.4, [[P("Une contrainte assumée : démontrer la faisabilité sur une infrastructure bancaire réelle, sans accélérateur.",
                                  12.5, SLATE, False, BODY_F, True)]], align=PP_ALIGN.CENTER)

# ---- Défense en profondeur ----
s = slide(); header(s, "04 · Réalisation", "Défense en profondeur — cinq scanners et un LLM", 14)
gif_box(s, "scanners.gif", 0.7, 1.8, 7.4, 5.0, fallback="sast_tools_overview.png", caption="Cinq outils SAST en parallèle — OWASP Top 10 2025", label="scanners en parallèle")
tools = [("Trivy","CVE et images"),("Gitleaks","secrets"),("Semgrep","OWASP code"),
         ("Checkov","IaC"),("OSV-Scanner","dépendances")]
yy=2.05
for t,d in tools:
    rect(s, 8.4, yy, 1.85, 0.42, fill=TEAL_D, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 8.4, yy+0.04, 1.85, 0.36, [[P(t, 12, WHITE, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, 10.4, yy+0.06, 2.3, 0.4, [[P(d, 12, BODY, False, BODY_F)]]); yy += 0.55
card_text(s, 8.4, 5.05, 4.25, 1.55, [
    [P("Complémentarité", 14, VIOLET_D, True, HEAD)],
    [P("le LLM a détecté des secrets en dur que Gitleaks n'avait pas signalés — scanners et IA se renforcent.",
       12, BODY, False, BODY_F)]], fill=T_VIOLET, line=None)

# ---- Démonstration PR #18 (BLOCK animé) ----
s = slide(); header(s, "04 · Réalisation", "Démonstration — Pull Request #18", 15)
gif_box(s, "demo_pr18.gif", 0.7, 1.85, 7.0, 4.95,
        fallback="pr_comment-security-review_0.png",
        caption="Revue scrollée automatiquement jusqu'au verdict BLOCK")
txt(s, 8.0, 2.05, 4.6, 0.5, [[P("Le déroulement, sans intervention humaine", 15, INK, True, HEAD)]])
yy=2.65
for t in ["Réception du webhook","Cinq scanners et LLM","Commentaires inline publiés","Suggestions de correction"]:
    rect(s, 8.0, yy+0.04, 0.16,0.16, fill=TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, 8.3, yy-0.02, 4.5, 0.4, [[P(t, 13, BODY, False, BODY_F)]]); yy+=0.55
badge = rect(s, 8.0, 5.2, 4.6, 0.7, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
badge.shadow.inherit=False; add_shadow(badge)
tf=badge.text_frame; tf.word_wrap=True
r=tf.paragraphs[0].add_run(); r.text="VERDICT :  BLOCK"; r.font.size=Pt(20); r.font.bold=True
r.font.name=HEAD; r.font.color.rgb=WHITE; tf.paragraphs[0].alignment=PP_ALIGN.CENTER
tf.vertical_anchor=MSO_ANCHOR.MIDDLE
txt(s, 8.0, 6.0, 4.6, 0.4, [[P("→ statut de commit en échec, fusion bloquée", 12, SLATE, False, BODY_F, True)]])
ANIM_BUILDS.append((s, [badge]))

# ---- Bénéfices et résultats ----
s = slide(); header(s, "04 · Réalisation", "Bénéfices et résultats", 16)
rows=[("Critère","Revue manuelle (AS-IS)","Agent IA (TO-BE)"),
      ("Délai","jusqu'à 24 h (selon le relecteur)","15–25 min"),
      ("Couverture","variable","OWASP Top 10 systématique"),
      ("Gate CI/CD","aucun","APPROVE / CHANGES / BLOCK"),
      ("Confidentialité","—","100 % local (BCT)")]
tb=s.shapes.add_table(len(rows),3,Inches(0.7),Inches(1.95),Inches(7.6),Inches(3.7)).table
tb.columns[0].width=Inches(2.0); tb.columns[1].width=Inches(2.9); tb.columns[2].width=Inches(2.7)
for r in range(len(rows)):
    for c in range(3):
        cell=tb.cell(r,c); cell.margin_left=Inches(0.1); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
        run=cell.text_frame.paragraphs[0].add_run(); run.text=rows[r][c]
        run.font.name=HEAD if (r==0 or c==0) else BODY_F; run.font.size=Pt(12 if r else 12.5)
        if r==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=(GREEN if c==2 else INK); run.font.color.rgb=WHITE; run.font.bold=True
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb=WHITE if r%2 else PANEL; run.font.color.rgb=BODY
            if c==0: run.font.color.rgb=INK; run.font.bold=True
            if c==2: run.font.color.rgb=RGBColor(0x12,0x7A,0x46); run.font.bold=True
card_text(s, 8.6, 1.95, 4.03, 3.7, [
    [P("5 / 5", 54, GREEN, True, HEAD)],
    [P("vulnérabilités plantées", 14, INK, True, HEAD)],
    [P("détectées sur la PR #18", 12, SLATE, False, BODY_F)],
    [P("", 6, BODY)],
    [P("+ corpus de huit revues réelles", 12, BODY, False, BODY_F)],
], anchor=MSO_ANCHOR.MIDDLE)
for pp in s.shapes[-1].text_frame.paragraphs: pp.alignment=PP_ALIGN.CENTER

# ---- Observabilité ----
s = slide(); header(s, "04 · Réalisation", "Supervision complète du VPS", 17)
fit_image(s, img("monitoring_stack.png"), 0.7, 1.85, 6.2, 4.6, caption="Stack d'observabilité")
gif_box(s, "grafana_live.gif", 7.2, 1.9, 5.4, 2.35, fallback="grafana_pr_reviews_dashboard.png", label="dashboard en direct")
fit_image(s, img("grafana_vps_host_dashboard.png"), 7.2, 4.45, 5.4, 1.9, card=True)
txt(s, 7.2, 6.45, 5.4, 0.4, [[P("Prometheus · VictoriaMetrics · Grafana · AlertManager  —  ", 11, SLATE, False, BODY_F),
                              P("28 métriques · 15 alertes · 3 dashboards", 11, TEAL_D, True, HEAD)]],
    align=PP_ALIGN.CENTER)

# ---- Assistant conversationnel & autonomie ----
s = slide(); header(s, "04 · Réalisation", "Assistant conversationnel et opérations autonomes", 18)
gif_box(s, "chat_demo.gif", 0.7, 1.85, 6.4, 4.9, fallback="inter_react.png", caption="Boucle ReAct — dix-neuf outils de supervision", label="chat répondant en direct")
txt(s, 7.4, 2.05, 5.2, 0.6, [[P("« Quel est l'usage CPU ? »", 16, VIOLET_D, True, HEAD)]])
txt(s, 7.4, 2.65, 5.2, 0.9, [[P("L'agent invoque un outil réel ", 13, BODY, False, BODY_F),
    P("(vps_status)", 12, INK, True, MONO), P(" et répond avec la valeur mesurée — jamais inventée.", 13, BODY, False, BODY_F)]])
txt(s, 7.4, 3.75, 5.2, 0.5, [[P("Et il s'auto-exploite :", 14, INK, True, HEAD)]])
for i,t in enumerate(["Scheduler autonome (tâches asyncio)","Health digest Slack quotidien (09:00 UTC)",
                      "Gardien disque auto-réparateur à 90 %"]):
    rect(s, 7.4, 4.3+i*0.62+0.03, 0.16,0.16, fill=GOLD, shape=MSO_SHAPE.OVAL)
    txt(s, 7.7, 4.3+i*0.62-0.03, 4.9, 0.5, [[P(t, 13, BODY, False, BODY_F)]])

# ---- Détection autonome des pannes ----
s = slide(); header(s, "04 · Réalisation", "Détection autonome de ses propres pannes", 19)
blob(s, 10.6, 4.0, 2.3, 2.3, T_VIOLET)
card_text(s, 0.7, 2.1, 5.6, 3.1, [
    [P("0 → 51", 56, VIOLET_D, True, HEAD)],
    [P("tokens", 16, INK, True, HEAD)], [P("", 8, BODY)],
    [P("La télémétrie A/B a révélé un bug du modèle Phi-4", 13, BODY, False, BODY_F)],
    [P("qui renvoyait zéro token, sans aucune erreur journalisée.", 13, BODY, False, BODY_F)],
], fill=PANEL, line=None, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 6.7, 2.4, 5.9, 0.6, [[P("La preuve d'un système qui se surveille lui-même", 17, INK, True, HEAD)]])
for i,t in enumerate(["Le système supervise sa propre infrastructure",
                      "Une anomalie invisible aux journaux est captée par sa métrique",
                      "Diagnostic, correctif puis vérification (0 → 51 tokens)"]):
    rect(s, 6.7, 3.3+i*0.8+0.03, 0.16,0.16, fill=VIOLET, shape=MSO_SHAPE.OVAL)
    txt(s, 7.0, 3.25+i*0.8, 5.5, 0.7, [[P(t, 13, BODY, False, BODY_F)]])

# ---- Tests et validation ----
s = slide(); header(s, "04 · Réalisation", "Tests et validation de la plateforme")
tests = [("HMAC-SHA256", "signature invalide rejetée en 403 ; signature valide acceptée en 202", TEAL),
         ("Déduplication Redis", "deux webhooks identiques ne déclenchent qu'un seul pipeline (SET NX)", VIOLET),
         ("Circuit breaker LLM", "Ollama indisponible : repli sur la classification par extensions, sans blocage", GOLD),
         ("Gardien disque", "au-delà de 90 % : nettoyage Docker automatique et notification Slack", INK)]
for i,(t,d,c) in enumerate(tests):
    col = i % 2; row = i // 2
    x = 0.7 + col*6.08; y = 1.95 + row*1.7
    sp = card_text(s, x, y, 5.85, 1.5, [
        [P("✓  ", 15, GREEN, True, HEAD), P(t, 14.5, INK, True, HEAD)],
        [P(d, 12, BODY, False, BODY_F)],
    ])
    rect(s, x, y, 0.1, 1.5, fill=c)
badge = rect(s, 0.7, 5.6, 11.93, 0.85, fill=RGBColor(0xEA,0xF7,0xF0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 0.7, 5.78, 11.93, 0.5, [[P("Checklist de validation finale : ", 14, INK, True, HEAD),
    P("quatorze vérifications de bout en bout, toutes validées avant la livraison.", 14, GREEN, True, HEAD)]],
    align=PP_ALIGN.CENTER)

# ======================= CHAPITRE 5 — CONCLUSION =======================
s = slide(); section_divider(s, 4, "Limites, perspectives et bilan du projet")

# ---- Limites et perspectives ----
s = slide(); header(s, "05 · Conclusion", "Limites assumées et perspectives", 20)
card_text(s, 0.7, 2.0, 5.85, 4.3, [
    [P("Limites assumées", 17, AMBER, True, HEAD)], [P("", 6, BODY)],
    [P("•  Inférence CPU : 15–25 min/PR — le prix de la confidentialité", 13, BODY, False, BODY_F)], [P("", 4, BODY)],
    [P("•  Corpus d'évaluation restreint (jeu contrôlé + 8 revues réelles)", 13, BODY, False, BODY_F)], [P("", 4, BODY)],
    [P("•  Diff = entrée non fiable → findings scanners déterministes,", 13, BODY, False, BODY_F)],
    [P("    ancrage renforcé du verdict prévu (défense en profondeur)", 13, BODY, False, BODY_F)],
], fill=PANEL, line=None)
card_text(s, 6.78, 2.0, 5.85, 4.3, [
    [P("Perspectives", 17, GREEN, True, HEAD)], [P("", 6, BODY)],
    [P("•  Migration GPU : 30–50 tok/s → 2–5 min/PR (~3000 €)", 13, BODY, False, BODY_F)], [P("", 4, BODY)],
    [P("•  Passage à l'échelle de tous les dépôts BTE", 13, BODY, False, BODY_F)], [P("", 4, BODY)],
    [P("•  Tableau de bord de risque pour le RSSI", 13, BODY, False, BODY_F)],
], fill=T_TEAL, line=None)

# ---- Conclusion & remerciements ----
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE); deco(s, "divider")
chevron(s, 0.9, 1.2, 0.3, VIOLET, n=3, gap=0.27)
txt(s, 0.9, 1.7, 11.5, 0.4, [[P("CONCLUSION", 14, TEAL_D, True, BODY_F)]])
txt(s, 0.9, 2.15, 11.6, 1.8, [
    [P("Un agent IA agentique, déployé en production à la BTE,", 25, INK, False, HEAD)],
    [P("qui sécurise chaque Pull Request ", 25, INK, False, HEAD),
     P("sans qu'aucun code ne quitte la banque.", 25, VIOLET_D, False, HEAD)],
])
for i,(b,l,c) in enumerate([("24 h → 15 min","délai",VIOLET_D),("5 / 5","détection",GREEN),("100 % local","conformité BCT",TEAL_D)]):
    x=0.9+i*4.0
    card_text(s, x, 4.2, 3.7, 1.3, [[P(b,22,c,True,HEAD)],[P(l,12,SLATE,False,BODY_F)]], anchor=MSO_ANCHOR.MIDDLE)
    for pp in s.shapes[-1].text_frame.paragraphs: pp.alignment=PP_ALIGN.CENTER
txt(s, 0.9, 5.9, 11.5, 0.6, [[P("Plateforme remise à la BTE   ·   Merci de votre attention", 17, GOLD, True, HEAD)]])

# ===================== animations =====================
for sl in prs.slides:
    add_fade_transition(sl)
for sl, shapes in ANIM_BUILDS:
    add_fade_in_build(sl, shapes)

# ===================== discours (notes orateur) =====================
SPEECH = [
# 1 — Titre (≈ 30 s)
"(≈ 30 s) Monsieur le Président du jury, Madame, Messieurs les membres du jury, bonjour. "
"Je m'appelle Ghaith Ferchichi et j'ai l'honneur de vous présenter mon projet de fin d'études, "
"réalisé au sein de la Banque de Tunisie et des Émirats, sous l'encadrement de Monsieur Kamel "
"Kaouech, côté banque, et de Madame Ghayet El Mouna Zhioua, côté ISI. Il s'intitule BTE Security "
"AI Agent : un agent d'intelligence artificielle pour la revue automatisée de sécurité du code.",
# 2 — Sommaire (≈ 10 s)
"(≈ 10 s) Mon exposé suit la structure du rapport : le contexte général, l'état de l'art et les "
"choix technologiques, la conception de l'agent, sa réalisation et ses résultats, puis la "
"conclusion et les perspectives.",
# 3 — Introduction (≈ 35 s)
"(≈ 35 s) Partons de trois scénarios concrets : une injection SQL fusionnée dans une API de "
"paiement, un secret d'authentification commité par mégarde dans un dépôt, une image Docker bâtie "
"sur une bibliothèque exposée à un CVE critique. Chacune de ces failles peut compromettre le "
"système d'information de la banque. Or, à la BTE, la revue de sécurité du code est entièrement "
"manuelle : elle peut atteindre vingt-quatre heures, et sa qualité dépend de la disponibilité et "
"de l'expertise du relecteur du jour. C'est ce constat qui motive ce travail.",
# 4 — Divider 1 (≈ 5 s)
"(≈ 5 s) Commençons par le contexte général.",
# 5 — Organisme & AS-IS (≈ 30 s)
"(≈ 30 s) La BTE est née en 1982 d'une convention entre l'État tunisien et l'Abu Dhabi Investment "
"Authority. Mon stage s'est déroulé à la Direction Centrale de l'Informatique et de "
"l'Organisation, auprès des équipes de sécurité opérationnelle. Le processus existant, à gauche, "
"repose sur GitHub : à l'ouverture d'une Pull Request, un développeur senior ou un référent "
"sécurité examine le diff à la main, commente, puis déclenche la fusion. Aucune analyse statique, "
"aucun scanner, aucun gate de pipeline : tout repose sur l'examen humain.",
# 6 — Problématique (≈ 30 s)
"(≈ 30 s) Quatre limites en découlent. Le délai, d'abord : jusqu'à vingt-quatre heures, pendant "
"lesquelles la vulnérabilité reste exposée. L'absence de blocage technique, ensuite : une Pull "
"Request porteuse d'une faille critique peut être fusionnée sans obstacle. Une qualité de revue "
"qui varie d'un relecteur à l'autre. Et deux angles morts structurels : les secrets commités par "
"accident et les dépendances vulnérables, deux catégories pour lesquelles l'œil humain n'est pas "
"le bon outil.",
# 7 — Objectifs et contribution (builds : 3 clics) (≈ 40 s)
"(≈ 40 s) Notre contribution tient en une phrase : un agent d'intelligence artificielle autonome "
"qui revoit chaque Pull Request selon la grille OWASP Top 10, publie ses correctifs directement "
"sur GitHub, et s'exécute intégralement en local — aucun code ne quitte la banque. Trois chiffres "
"résument le résultat. [clic] Le délai de revue passe de vingt-quatre heures à une quinzaine de "
"minutes. [clic] Sur la validation contrôlée, cinq vulnérabilités sur cinq sont détectées, sans "
"faux positif dans le fichier modifié. [clic] Et cent pour cent de l'inférence reste sur le VPS "
"de la banque, conformément aux exigences de la Banque Centrale de Tunisie.",
# 8 — Solution TO-BE (≈ 30 s)
"(≈ 30 s) Voici le processus cible. GitHub émet un webhook signé HMAC-SHA256 ; l'agent vérifie la "
"signature, élimine les doublons, classifie la Pull Request en une trentaine de secondes, exécute "
"en parallèle les scanners pertinents, puis consolide leurs résultats avec un modèle de langage "
"local. Il publie sur GitHub la revue complète : score de risque, commentaires ligne par ligne et "
"verdict — APPROVE, REQUEST_CHANGES ou BLOCK — posé comme statut de commit. C'est ce statut qui "
"conditionne la fusion, laquelle reste déclenchée par un humain.",
# 9 — Méthodologie (≈ 30 s)
"(≈ 30 s) Le projet a été conduit selon la méthode Scrumban, retenue après comparaison avec "
"Scrum, Kanban et XP : elle combine la cadence et les jalons de Scrum, utiles à l'encadrement, "
"avec la discipline de flux de Kanban — un travail en cours limité à une seule carte. Dix sprints "
"de deux semaines sur cinq mois, six jalons, et une étiquette Expedite qui permettait à un "
"incident de production de passer devant tout le reste, sans replanifier le sprint en cours.",
# 10 — Divider 2 (≈ 5 s)
"(≈ 5 s) Passons à l'état de l'art et aux choix technologiques.",
# 11 — DevSecOps (≈ 25 s)
"(≈ 25 s) Le cadre conceptuel est le DevSecOps et son principe de shift-left : déplacer la "
"détection des vulnérabilités au plus tôt du cycle, idéalement dès la Pull Request, là où une "
"faille ne coûte presque rien à corriger comparée à la même faille découverte en production. "
"Notre apport matérialise ce principe : le verdict de l'agent devient un gate technique, qui "
"conditionne réellement la fusion.",
# 12 — Besoins (≈ 25 s)
"(≈ 25 s) Cinq besoins fonctionnels : détecter les vulnérabilités à chaque Pull Request, produire "
"la revue OWASP avec score et verdict, publier les commentaires sur GitHub, persister chaque "
"revue en base de connaissances, et offrir un assistant d'exploitation en langage naturel. Côté "
"non fonctionnel, la contrainte dominante est la confidentialité — l'inférence doit rester locale "
"—, complétée par la performance, la résilience par circuit breaker et l'observabilité complète.",
# 13 — Benchmark (≈ 45 s)
"(≈ 45 s) Le choix des modèles ne repose pas sur la littérature, mais sur un banc d'essai mené "
"sur le matériel cible, avec le prompt système complet. Deux candidats sont éliminés : llama 3.2 "
"trois milliards, dont la fenêtre de contexte est saturée par le seul prompt système — sa "
"précision tombe à zéro —, et granite, dont le schéma d'appel d'outils est incompatible. Restent "
"qwen2.5-coder 7B et 14B, à égalité de précision à quatre-vingts pour cent. D'où une architecture "
"à deux modèles : le 7B classifie en trente secondes, le 14B mène l'analyse approfondie. Un "
"benchmark croisé a par ailleurs confirmé le moteur Ollama : vingt-deux pour cent de débit de "
"plus que LocalAI, sur un fichier de poids strictement identique.",
# 14 — Technologies retenues (≈ 30 s)
"(≈ 30 s) Trois briques portent l'ensemble. Cinq scanners spécialisés : Trivy pour les CVE, "
"Gitleaks pour les secrets, Semgrep pour le code applicatif, Checkov pour l'infrastructure, "
"OSV-Scanner pour les dépendances. Deux modèles servis localement par Ollama. Et l'orchestrateur "
"LangGraph, préféré à Prefect et Celery pour sa persistance d'état native dans PostgreSQL : un "
"pipeline de vingt minutes survit au redémarrage de son conteneur. Détail qui compte : les "
"sorties des scanners sont nettoyées avant transmission au modèle — cinquante-deux pour cent de "
"tokens en moins — pour préserver la fenêtre de contexte.",
# 15 — Positionnement (≈ 30 s)
"(≈ 30 s) Pourquoi ne pas avoir adopté un outil du marché ? Parce qu'aucun ne combine "
"raisonnement par modèle de langage et inférence entièrement locale. CodeQL, Snyk et Copilot "
"Autofix transmettent le code à une infrastructure tierce — rédhibitoire en contexte bancaire. "
"SonarQube est auto-hébergeable, mais purement statique, sans explication en langage naturel. "
"Notre contribution est l'intégration de ces deux mondes, sans qu'une ligne de code ne sorte du "
"VPS.",
# 16 — Divider 3 (≈ 5 s)
"(≈ 5 s) Venons-en à la conception.",
# 17 — Architecture (≈ 25 s)
"(≈ 25 s) La plateforme compte onze conteneurs Docker en quatre couches. L'entrée : nginx, seul "
"point exposé, qui authentifie et transmet les webhooks. La couche IA : l'agent FastAPI, qui "
"orchestre le graphe LangGraph et les scanners, et Ollama, accessible uniquement depuis le réseau "
"Docker interne. Les données : PostgreSQL pour la base de connaissances et les checkpoints, Redis "
"pour la déduplication et le cache. Et l'observabilité, sur laquelle je reviendrai.",
# 18 — Manifeste (builds : 6 clics) (≈ 35 s)
"(≈ 35 s) Ce point fonde le titre du projet : il s'agit d'un agent, et non d'un simple appel à un "
"modèle. Six capacités le démontrent. [clics] Il perçoit son environnement — webhooks, métriques, "
"alertes. Il raisonne : il classifie chaque Pull Request et route son graphe d'état en "
"conséquence. Il agit : il commente, il bloque, il escalade vers Slack. Il persiste : son état "
"est sauvegardé après chaque nœud, et il reprend exactement où il s'était arrêté. Il manipule "
"dix-neuf outils réels. Et il s'auto-exploite : gardien disque et bilan de santé quotidien.",
# 19 — Pipeline (builds : 7 clics) (≈ 50 s)
"(≈ 50 s) Le cœur du système est un graphe d'état à neuf nœuds. Suivons une Pull Request. [clics] "
"Le webhook est vérifié et dédupliqué ; le dépôt est cloné et le diff généré localement avec "
"quinze lignes de contexte. Le modèle 7B classifie en cinq catégories. Le routage adapte "
"l'analyse : une PR de documentation saute les scanners, un Dockerfile déclenche le scan d'image "
"complet. Les scanners s'exécutent en parallèle, chacun isolé dans sa coroutine — si l'un échoue, "
"le pipeline continue. Le 14B mène ensuite l'analyse OWASP sur un diff annoté de numéros de "
"lignes ; le verdict est rendu, avec escalade Slack si le risque est élevé, puis la revue est "
"publiée. Une optimisation a changé l'échelle : la fusion des deux appels du 14B en un seul a "
"fait passer le pipeline de trente-cinquante minutes à quinze-vingt-cinq.",
# 20 — Fiabilité (≈ 35 s)
"(≈ 35 s) Reste la question légitime : peut-on faire confiance à un modèle de langage ? Six "
"couches de protection répondent. Trois agissent sur le modèle : température proche de zéro, "
"fenêtre de contexte dimensionnée, génération plafonnée. Trois relèvent du contrôle applicatif : "
"une garde qui force l'appel d'outil pour toute donnée en temps réel, des règles "
"anti-hallucination dans le prompt, et l'interdiction de citer une valeur absente des "
"observations. Pour la revue de code s'y ajoute un parser dédié, qui écarte tout commentaire "
"visant une ligne inexistante du diff. Résultat : aucune erreur de ligne publiée sur GitHub.",
# 21 — Divider 4 (≈ 5 s)
"(≈ 5 s) Passons à la réalisation et aux résultats.",
# 22 — Environnement (≈ 20 s)
"(≈ 20 s) Le terrain d'exécution est volontairement contraint : un VPS de production de douze "
"cœurs Haswell, quarante-cinq giga-octets de mémoire, sans GPU — c'est le prix de la "
"confidentialité, et nous l'assumons. Cette contrainte a imposé un vrai tuning : flash attention "
"et quantisation du cache KV pour tenir le modèle 14B en mémoire. La pile est intégralement open "
"source.",
# 23 — Défense en profondeur (≈ 30 s)
"(≈ 30 s) La détection suit une logique de défense en profondeur : chaque scanner couvre une "
"classe de risque, et le modèle consolide l'ensemble en une revue unique. Les deux niveaux se "
"complètent réellement : sur la démonstration, c'est le modèle qui a détecté des secrets en dur "
"que Gitleaks n'avait pas signalés. Et les findings des scanners sont déterministes : ils sont "
"persistés en base quel que soit le texte produit par le modèle.",
# 24 — Démonstration (GIF + build BLOCK) (≈ 50 s)
"(≈ 50 s) Voici la démonstration, sur la Pull Request numéro dix-huit : un module bancaire en PHP "
"contenant cinq vulnérabilités types — injection SQL, secrets en dur, path traversal, injection "
"de commande et hachage MD5. À l'écran, la revue publiée par l'agent défile : l'analyse de chaque "
"faille, puis les recommandations. [clic] Le verdict tombe : BLOCK. Cinq commentaires inline, "
"chacun avec une suggestion de correction applicable en un clic depuis GitHub, et le statut de "
"commit passe en échec : la fusion est techniquement bloquée tant que les corrections ne sont pas "
"poussées. Durée totale pour cette revue : une quinzaine de minutes, notification Slack comprise.",
# 25 — Bénéfices & résultats (≈ 25 s)
"(≈ 25 s) Le tableau résume le passage du manuel à l'agent : un délai divisé par près de cent, "
"une couverture OWASP systématique, un gate effectif là où rien ne bloquait, une traçabilité "
"complète en base. Sur le jeu contrôlé, cinq sur cinq, sans faux positif dans le fichier modifié. "
"Au-delà de la démonstration, huit revues réelles sont persistées en base, et leurs verdicts "
"suivent le risque.",
# 26 — Observabilité (≈ 35 s)
"(≈ 35 s) L'agent voit l'intégralité du VPS : vingt-huit métriques personnalisées, quinze règles "
"d'alerte, trois tableaux de bord Grafana. Les dernières métriques interrogent directement le "
"socket Docker — après l'échec de cAdvisor sur notre version de Docker — et comparent les deux "
"moteurs d'inférence. Cette supervision n'est pas décorative : elle a révélé un AlertManager "
"silencieusement hors service depuis le déploiement initial, et une panne de VictoriaMetrics "
"restée neuf jours invisible. Deux anomalies détectées, puis corrigées en sprint.",
# 27 — Chat & autonomie (≈ 25 s)
"(≈ 25 s) Un assistant conversationnel doté de dix-neuf outils permet d'interroger "
"l'infrastructure en langage naturel, jusqu'à la base de connaissances sécurité, sans écrire de "
"SQL. Chaque question sur une donnée réelle déclenche l'appel d'outil correspondant : la réponse "
"cite la valeur mesurée, jamais inventée. Et l'agent s'auto-exploite : gardien disque toutes les "
"trente minutes, nettoyage automatique à quatre-vingt-dix pour cent, bilan de santé publié chaque "
"matin sur Slack.",
# 28 — Phi-4 (≈ 35 s)
"(≈ 35 s) L'épisode le plus révélateur du projet : la télémétrie a mis en évidence un bug du "
"modèle Phi-4, qui renvoyait zéro token alors que la requête se terminait sans aucune erreur "
"journalisée. La cause : un template de conversation mal formé, qui présentait au modèle une "
"conversation déjà close. Le correctif a été aligné sur la configuration officielle du modèle, "
"puis vérifié : de zéro à cinquante et un tokens. Le système a détecté une panne invisible aux "
"journaux — c'est la meilleure preuve qu'il se surveille lui-même.",
# 29 — Tests et validation (≈ 20 s)
"(≈ 20 s) La validation finale a couvert les mécanismes critiques : signature invalide rejetée "
"avant tout traitement, webhook dupliqué ignoré, circuit breaker avec classification de repli si "
"le modèle est indisponible, gardien disque vérifié en conditions réelles. Au total, quatorze "
"vérifications de bout en bout, toutes validées avant la livraison.",
# 30 — Divider 5 (≈ 5 s)
"(≈ 5 s) J'en viens à la conclusion.",
# 31 — Limites assumées & perspectives (≈ 40 s)
"(≈ 40 s) Deux limites assumées et un point de vigilance. L'inférence sur CPU borne la revue à "
"quinze-vingt-cinq minutes : c'est le prix de la confidentialité. Le corpus d'évaluation reste "
"restreint — un jeu contrôlé et huit revues réelles. Et le diff analysé est par nature une entrée "
"non fiable : l'architecture y oppose déjà les findings déterministes des scanners, qu'aucun "
"texte ne peut effacer, et le plan de durcissement prévoit d'ancrer davantage le verdict sur ces "
"findings. Trois perspectives : la migration GPU — environ trois mille euros, dans l'enveloppe de "
"la DCIO — ramènerait chaque revue à deux à cinq minutes, sans toucher ni à l'architecture ni à "
"la confidentialité ; le passage à l'échelle de l'ensemble des dépôts de la banque ; et un "
"tableau de bord de risque alimenté par les revues accumulées, destiné au RSSI.",
# 32 — Conclusion (≈ 30 s)
"(≈ 30 s) En conclusion : parti d'un VPS vierge, ce projet livre un agent d'intelligence "
"artificielle complet, déployé en production à la BTE, qui sécurise chaque Pull Request sans "
"qu'aucun code ne quitte la banque. La plateforme est remise à la banque et constitue un socle "
"concret pour sa transformation DevSecOps. Je vous remercie de votre attention et me tiens à "
"votre disposition pour répondre à vos questions.",
]
_slides = list(prs.slides)
if len(SPEECH) != len(_slides):
    print(f"ATTENTION : {len(SPEECH)} notes pour {len(_slides)} slides")
for sl, note_text in zip(_slides, SPEECH):
    sl.notes_slide.notes_text_frame.text = note_text

prs.save(OUT)
print("OK ->", OUT, "| slides:", len(prs.slides._sldIdLst), "| builds:", len(ANIM_BUILDS), "| notes:", len(SPEECH))
