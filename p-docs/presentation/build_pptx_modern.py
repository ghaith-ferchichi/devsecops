#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Soutenance — BTE Security AI Agent.
Thème « Modern Light Tech » (blanc + violet IA + teal + or, sans bold,
images/icônes, animations fluides). Réf : presentation/pptx-plan.md §3/§9/§11.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from PIL import Image

_HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.normpath(os.path.join(_HERE, "..", "MASTER_PFE", "img"))
OUT = os.path.join(_HERE, "BTE_Security_AI_Agent_Soutenance.pptx")
GIFS = os.path.join(_HERE, "gifs")   # GIFs animés (auto-loop en mode Présentation)

# ---------- Palette : Modern Light Tech ----------
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
PANEL   = RGBColor(0xF2, 0xF4, 0xF7)
INK     = RGBColor(0x0B, 0x25, 0x45)
BODY    = RGBColor(0x33, 0x40, 0x5C)
VIOLET  = RGBColor(0x7C, 0x5C, 0xFC)   # AI fill
VIOLET_D= RGBColor(0x6D, 0x28, 0xD9)   # AI text
TEAL    = RGBColor(0x2B, 0xA6, 0xAE)   # tech fill
TEAL_D  = RGBColor(0x15, 0x7A, 0x82)   # tech text/lines
GOLD    = RGBColor(0xE0, 0xA1, 0x00)   # BTE fill
SLATE   = RGBColor(0x5B, 0x6B, 0x7B)
HAIR    = RGBColor(0xE1, 0xE7, 0xED)
GREEN   = RGBColor(0x1E, 0x9E, 0x5A)
AMBER   = RGBColor(0xE0, 0x86, 0x00)
RED     = RGBColor(0xD7, 0x26, 0x3D)
# light tints for decorative blobs
T_VIOLET= RGBColor(0xEF, 0xEA, 0xFE)
T_TEAL  = RGBColor(0xE6, 0xF6, 0xF7)
T_GOLD  = RGBColor(0xFBF1D8 >> 16 & 255, 0xFBF1D8 >> 8 & 255, 0xFBF1D8 & 255)

HEAD = "Segoe UI Semibold"
BODY_F = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
ANIM_BUILDS = []   # list of (slide, [shapes]) for fade-in entrance builds

# ===================== low-level helpers =====================
def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow: add_shadow(sp)
    return sp

def add_shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = spPr.makeelement(qn('a:outerShdw'),
        {'blurRad':'80000','dist':'30000','dir':'5400000','rotWithShape':'0'})
    clr = spPr.makeelement(qn('a:srgbClr'), {'val':'0B2545'})
    alpha = spPr.makeelement(qn('a:alpha'), {'val':'18000'})
    clr.append(alpha); sh.append(clr); el.append(sh); spPr.append(el)

def round_shape(sp, radius_in, w_in, h_in):
    """Force prstGeom=roundRect with an absolute corner radius (inches)."""
    spPr = sp._element.spPr
    frac = max(0.0, min(0.5, radius_in / min(w_in, h_in)))
    val = str(int(frac * 100000))
    geom = spPr.find(qn('a:prstGeom'))
    if geom is None:
        geom = parse_xml('<a:prstGeom xmlns:a="http://schemas.openxmlformats.org/'
                         'drawingml/2006/main" prst="roundRect"><a:avLst/></a:prstGeom>')
        xfrm = spPr.find(qn('a:xfrm'))
        xfrm.addnext(geom)
    geom.set('prst', 'roundRect')
    avLst = geom.find(qn('a:avLst'))
    if avLst is None:
        avLst = geom.makeelement(qn('a:avLst'), {}); geom.append(avLst)
    for gd in list(avLst): avLst.remove(gd)
    avLst.append(avLst.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': f'val {val}'}))

def set_gradient(sp, c1, c2, angle_deg=0):
    spPr = sp._element.spPr
    for t in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        e = spPr.find(qn(t))
        if e is not None: spPr.remove(e)
    ang = int(angle_deg*60000)
    xml = ('<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
           '<a:gsLst>'
           f'<a:gs pos="0"><a:srgbClr val="{str(c1)}"/></a:gs>'
           f'<a:gs pos="100000"><a:srgbClr val="{str(c2)}"/></a:gs>'
           '</a:gsLst>'
           f'<a:lin ang="{ang}" scaled="1"/></a:gradFill>')
    grad = parse_xml(xml)
    ln = spPr.find(qn('a:ln'))
    if ln is not None: ln.addprevious(grad)
    else: spPr.append(grad)

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, sp_after=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        for (t, sz, col, bold, fnt, *rest) in para:
            ital = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = ital
            r.font.name = fnt; r.font.color.rgb = col
    return tb

def P(t, sz, col=BODY, bold=False, fnt=BODY_F, ital=False):
    return (t, sz, col, bold, fnt, ital)

def card_text(s, x, y, w, h, runs, fill=WHITE, line=HAIR, accent=None,
              shadow=True, anchor=MSO_ANCHOR.TOP, pad=0.28):
    """A single rounded-rect SHAPE whose own text frame holds `runs`
    (so it animates as ONE shape). Optional left accent bar drawn separately."""
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    if shadow: add_shadow(sp)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=Inches(pad); tf.margin_right=Inches(pad)
    tf.margin_top=Inches(0.18); tf.margin_bottom=Inches(0.14)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(3)
        for (t, sz, col, bold, fnt, *rest) in para:
            ital = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = ital
            r.font.name = fnt; r.font.color.rgb = col
    return sp

def fit_image(s, path, bx, by, bw, bh, card=True, caption=None):
    iw, ih = Image.open(path).size
    ar = iw/ih; bar = bw/bh
    if ar > bar: w = bw; h = bw/ar
    else: h = bh; w = bh*ar
    x = bx + (bw-w)/2; y = by + (bh-h)/2
    pad = 0.10
    img_r = 0.09  # corner radius of the image (inches)
    if card:
        c = rect(s, x-pad, y-pad, w+2*pad, h+2*pad, fill=WHITE, line=HAIR,
             line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        # card radius = image radius + pad → uniform rounded frame, no corner overhang
        round_shape(c, img_r+pad, w+2*pad, h+2*pad)
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    round_shape(pic, img_r, w, h)
    if caption:
        txt(s, bx, by+bh+0.05, bw, 0.3, [[P(caption, 10, SLATE, False, BODY_F, True)]],
            align=PP_ALIGN.CENTER)
    return (x, y, w, h)

def img(name):
    for sub in ("chapter_1","chapter_2","chapter_3","chapter_4","template"):
        p = os.path.join(IMG, sub, name)
        if os.path.exists(p): return p
    raise FileNotFoundError(name)

RES = os.path.join(_HERE, "pres-ressources")
def res(name):
    """Captures réelles de la démonstration (pres-ressources/)."""
    p = os.path.join(RES, name)
    if os.path.exists(p): return p
    raise FileNotFoundError(name)

def rounded_pic(s, path, x, y, w, h, adj=9000):
    """Insère une image découpée en rectangle à coins arrondis (adj = rayon, 1/1000 %)."""
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    spPr = pic._element.spPr
    for g in spPr.findall(qn('a:prstGeom')):
        spPr.remove(g)
    geom = parse_xml(
        '<a:prstGeom xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'prst="roundRect"><a:avLst><a:gd name="adj" fmla="val {adj}"/></a:avLst></a:prstGeom>')
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None: xfrm.addnext(geom)
    else: spPr.insert(0, geom)
    return pic

def gif_box(s, gif_name, bx, by, bw, bh, fallback=None, caption=None, label=None):
    """GIF animé (auto-loop en mode Présentation) cadré dans une carte.
    À défaut : image statique de repli + badge « ▶ GIF à venir » ; sinon placeholder."""
    gp = os.path.join(GIFS, gif_name)
    if os.path.exists(gp):
        return fit_image(s, gp, bx, by, bw, bh, card=True, caption=caption)
    if fallback:
        return fit_image(s, img(fallback), bx, by, bw, bh, card=True, caption=caption)
    rect(s, bx, by, bw, bh, fill=PANEL, line=VIOLET, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, bx, by+bh/2-0.5, bw, 1.0, [[P("▶", 34, VIOLET, True, HEAD)],
        [P("GIF à venir — " + (label or gif_name), 12, SLATE, False, BODY_F)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return (bx, by, bw, bh)

# ===================== decorative + icon helpers =====================
def chevron(s, x, y, size, color, n=3, gap=0.18):
    for i in range(n):
        c = rect(s, x+i*gap, y, size, size, fill=color, shape=MSO_SHAPE.CHEVRON)
    return

def dotted_grid(s, x, y, cols, rows, gap=0.22, d=0.05, color=HAIR):
    for r in range(rows):
        for c in range(cols):
            rect(s, x+c*gap, y+r*gap, d, d, fill=color, shape=MSO_SHAPE.OVAL)

def blob(s, x, y, w, h, color, shape=MSO_SHAPE.OVAL):
    rect(s, x, y, w, h, fill=color, shape=shape)

def arcs(s, cx, cy, color=HAIR, n=3, r0=0.5, step=0.35):
    for i in range(n):
        r = r0 + i*step
        rect(s, cx-r, cy-r, 2*r, 2*r, fill=None, line=color, line_w=1.2, shape=MSO_SHAPE.OVAL)

def deco(s, kind="light"):
    """Subtle background decoration. Call right after the white canvas."""
    if kind == "light":
        blob(s, 11.9, -0.7, 1.9, 1.9, T_VIOLET)
        dotted_grid(s, 11.0, 0.35, 6, 3, color=HAIR)
        blob(s, -0.6, 6.0, 1.7, 1.7, T_TEAL)
    elif kind == "divider":
        blob(s, 11.4, -0.9, 2.6, 2.6, T_VIOLET)
        blob(s, -0.9, 5.6, 2.4, 2.4, T_TEAL)
        arcs(s, 12.7, 7.4, color=RGBColor(0xEDE9FB>>16&255,0xEDE9FB>>8&255,0xEDE9FB&255), n=3, r0=0.5)
    elif kind == "hero":
        dotted_grid(s, 0.7, 0.35, 6, 3, color=HAIR)
        blob(s, 11.9, 5.9, 1.9, 1.9, T_GOLD)

def icon_badge(s, x, y, d, label, color):
    """Circular badge with a short text glyph/number (geometric, swap later)."""
    rect(s, x, y, d, d, fill=color, shape=MSO_SHAPE.OVAL)
    txt(s, x, y+d*0.16, d, d*0.7, [[P(label, 18, WHITE, True, HEAD)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===================== animation injection =====================
def add_fade_transition(sld_obj, spd="med"):
    sld = sld_obj._element
    for t in sld.findall(qn('p:transition')): sld.remove(t)
    xml = ('<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
           f'spd="{spd}"><p:fade/></p:transition>')
    el = parse_xml(xml)
    anchor = sld.find(qn('p:clrMapOvr'))
    if anchor is None: anchor = sld.find(qn('p:cSld'))
    sld.insert(list(sld).index(anchor)+1, el)

def add_fade_in_build(sld_obj, shapes, dur=500):
    """Standard mainSeq click-to-reveal fade-in; one click per shape."""
    sld = sld_obj._element
    for t in sld.findall(qn('p:timing')): sld.remove(t)
    ids = [sp.shape_id for sp in shapes]
    cid = [2]
    def nid():
        cid[0]+=1; return cid[0]
    steps = []
    for spid in ids:
        a,b,c,d,e = nid(),nid(),nid(),nid(),nid()
        steps.append(
          f'<p:par><p:cTn id="{a}" fill="hold"><p:stCondLst><p:cond delay="indefinite"/>'
          f'</p:stCondLst><p:childTnLst>'
          f'<p:par><p:cTn id="{b}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
          f'<p:childTnLst>'
          f'<p:par><p:cTn id="{c}" presetID="10" presetClass="entr" presetSubtype="0" '
          f'fill="hold" nodeType="clickEffect"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
          f'<p:childTnLst>'
          f'<p:set><p:cBhvr><p:cTn id="{d}" dur="1" fill="hold"><p:stCondLst>'
          f'<p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
          f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>'
          f'<p:to><p:strVal val="visible"/></p:to></p:set>'
          f'<p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{e}" dur="{dur}"/>'
          f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>'
          f'</p:childTnLst></p:cTn></p:par>'
          f'</p:childTnLst></p:cTn></p:par>'
          f'</p:childTnLst></p:cTn></p:par>')
    blds = "".join(f'<p:bldP spid="{i}" grpId="0" animBg="1"/>' for i in ids)
    xml = ('<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
           '<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
           '<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
           '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
           + "".join(steps) +
           '</p:childTnLst></p:cTn>'
           '<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
           '<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
           '</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
           f'<p:bldLst>{blds}</p:bldLst></p:timing>')
    el = parse_xml(xml)
    anchor = sld.find(qn('p:transition'))
    if anchor is None: anchor = sld.find(qn('p:clrMapOvr'))
    if anchor is None: anchor = sld.find(qn('p:cSld'))
    sld.insert(list(sld).index(anchor)+1, el)

# ===================== composite components =====================
_PAGE = [0]   # numérotation automatique des slides de contenu
def header(s, kicker, title, num=True, deco_kind="light"):
    rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
    deco(s, deco_kind)
    if kicker:
        txt(s, 0.7, 0.5, 11.0, 0.35, [[P(kicker.upper(), 12, TEAL_D, True, BODY_F)]])
    txt(s, 0.7, 0.8, 12.0, 0.9, [[P(title, 28, INK, False, HEAD)]])
    rect(s, 0.72, 1.62, 0.5, 0.07, fill=VIOLET, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 1.30, 1.63, 1.0, 0.05, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if num:
        _PAGE[0] += 1
        footer(s, _PAGE[0])
    else:
        footer(s, None)

def footer(s, num=None):
    rect(s, 0.7, 7.04, 11.93, 0.012, fill=HAIR)
    txt(s, 0.7, 7.1, 8.0, 0.3,
        [[P("BTE Security AI Agent", 9, SLATE, False, BODY_F),
          P("  ·  Ghaith Ferchichi", 9, SLATE, False, BODY_F)]])
    if num is not None:
        chevron(s, 11.55, 7.12, 0.12, VIOLET, n=2, gap=0.11)
        txt(s, 11.9, 7.07, 0.7, 0.3, [[P(f"{num:02d}", 10, INK, True, HEAD)]], align=PP_ALIGN.RIGHT)

PLAN_NODES = ["Contexte", "État de l'art", "Conception", "Réalisation", "Conclusion"]
def plan_divider(s, active, y=3.5):
    n = len(PLAN_NODES); x0 = 1.2; x1 = 12.13; d = 0.62
    span = x1 - x0; gap = span/(n-1)
    rect(s, x0+d/2, y+d/2-0.01, span, 0.02, fill=HAIR)
    for i, lab in enumerate(PLAN_NODES):
        cx = x0 + i*gap
        on = (i == active)
        c = VIOLET if on else PANEL
        dd = d if on else d*0.82
        rect(s, cx+(d-dd)/2, y+(d-dd)/2, dd, dd, fill=c,
             line=(None if on else HAIR), line_w=1.2, shape=MSO_SHAPE.OVAL)
        txt(s, cx+(d-dd)/2, y+(d-dd)/2+dd*0.16, dd, dd*0.7,
            [[P(str(i+1), 16 if on else 13, WHITE if on else SLATE, True, HEAD)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, cx-0.6, y+d+0.12, d+1.2, 0.4,
            [[P(lab, 12 if on else 10.5, INK if on else SLATE, on, BODY_F)]],
            align=PP_ALIGN.CENTER)

def section_divider(s, active, subtitle):
    rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
    deco(s, "divider")
    num = active+1
    txt(s, 0.7, 0.9, 4.0, 1.6, [[P(f"{num:02d}", 78, RGBColor(0xE7,0xE3,0xF7), True, HEAD)]])
    txt(s, 0.72, 2.0, 11.0, 0.9, [[P(PLAN_NODES[active], 40, INK, False, HEAD)]])
    rect(s, 0.74, 2.92, 0.5, 0.08, fill=VIOLET, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 1.32, 2.93, 1.2, 0.06, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 0.74, 3.12, 11.0, 0.5, [[P(subtitle, 15, SLATE, False, BODY_F, True)]])
    plan_divider(s, active, y=4.7)
    footer(s, None)

# ===================== SLIDES =====================

# ---- 1. PAGE DE TITRE (mise en page institutionnelle ISI, centrée) ----
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
# fines bandes dégradé (identité « Modern Light Tech ») en haut et en bas
_t = rect(s, 0, 0, 13.333, 0.15, fill=VIOLET); set_gradient(_t, VIOLET, TEAL, 0)
_b = rect(s, 0, 7.35, 13.333, 0.15, fill=TEAL);  set_gradient(_b, TEAL, VIOLET, 0)
# logos institutionnels : BTE (gauche) · ISI + UTM (droite)
def _logo(nm, x, y, h):
    try: s.shapes.add_picture(img(nm), Inches(x), Inches(y), height=Inches(h))
    except Exception: pass
_logo("logo_BTE.png", 0.62, 0.60, 0.50)     # AR 5.93 -> ~2.97 large
_logo("LogoISI.png", 10.78, 0.52, 0.56)     # AR 2.00 -> ~1.12 large
_logo("Logo_UTM.png", 12.16, 0.42, 0.72)    # AR 1.06 -> ~0.76 large
# en-tête institutionnel, centré entre les logos
txt(s, 3.0, 0.46, 7.33, 1.65, [
    [P("République Tunisienne", 13, INK, True, HEAD)],
    [P("Ministère de l'Enseignement Supérieur et de la Recherche Scientifique", 11, BODY, False, BODY_F)],
    [P("Université de Tunis El Manar", 11, BODY, False, BODY_F)],
    [P("Institut Supérieur d'Informatique", 11.5, TEAL_D, True, HEAD)],
], align=PP_ALIGN.CENTER, sp_after=2)
# filet de séparation
rect(s, 1.6, 2.30, 10.13, 0.018, fill=HAIR)
# acte de soutenance + spécialité (centré)
txt(s, 1.0, 2.46, 11.33, 0.45, [[P("Soutenance de Stage de Fin d'Études", 18, TEAL_D, True, HEAD)]],
    align=PP_ALIGN.CENTER)
txt(s, 1.0, 2.92, 11.33, 0.7, [
    [P("En vue de l'obtention du Diplôme National de Mastère Professionnel", 12, INK, True, HEAD)],
    [P("Spécialité : Sécurité des Systèmes d'Informations et des Infrastructures", 11.5, SLATE, False, BODY_F)],
], align=PP_ALIGN.CENTER, sp_after=2)
# carte-titre centrée + barre d'accent dégradé (identité moderne)
_cw, _ch = 8.4, 1.66
_cx = (13.333 - _cw) / 2
_cy = 3.80
_card = rect(s, _cx, _cy, _cw, _ch, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
round_shape(_card, 0.16, _cw, _ch)
_bar = rect(s, _cx + 0.18, _cy, _cw - 0.36, 0.10, fill=VIOLET, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
set_gradient(_bar, VIOLET, TEAL, 0)
txt(s, _cx, _cy + 0.30, _cw, 0.7,
    [[P("BTE Security ", 34, INK, False, HEAD), P("AI Agent", 34, VIOLET_D, False, HEAD)]],
    align=PP_ALIGN.CENTER)
txt(s, _cx, _cy + 1.08, _cw, 0.45,
    [[P("Agent IA pour la revue automatisée de code", 14, BODY, False, BODY_F)]],
    align=PP_ALIGN.CENTER)
# bloc bas : encadrants (gauche) · réalisé par (droite) · année (centre)
txt(s, 0.9, 5.82, 5.6, 1.2, [
    [P("Encadrants", 12, VIOLET_D, True, HEAD)],
    [P("M. Kamel KAOUECH ", 11, INK, True, HEAD), P("— professionnel", 10.5, SLATE, False, BODY_F)],
    [P("Mme Ghayet El Mouna ZHIOUA ", 11, INK, True, HEAD), P("— académique", 10.5, SLATE, False, BODY_F)],
], sp_after=3)
txt(s, 6.8, 5.82, 5.63, 1.2, [
    [P("Réalisé par", 12, TEAL_D, True, HEAD)],
    [P("Ghaith FERCHICHI", 13, INK, True, HEAD)],
], align=PP_ALIGN.RIGHT, sp_after=3)
txt(s, 1.0, 6.97, 11.33, 0.35,
    [[P("Année universitaire 2025–2026", 10.5, SLATE, False, BODY_F, True)]],
    align=PP_ALIGN.CENTER)

# ---- 2. SOMMAIRE ----
s = slide(); header(s, "Plan", "Sommaire", None)
somm = [("1","Contexte général","Organisme d'accueil, étude de l'existant, problématique et méthodologie"),
        ("2","État de l'art et choix technologiques","Concepts DevSecOps, spécification des besoins et technologies retenues"),
        ("3","Conception","Architecture de l'agent — « un cerveau, et non un simple pipeline »"),
        ("4","Réalisation et résultats","Sécurité, démonstration, observabilité et validation"),
        ("5","Conclusion et perspectives","Limites, perspectives et bilan du projet")]
yy=1.95
for i,(n,t,d) in enumerate(somm):
    c = VIOLET if i % 2 == 0 else TEAL
    icon_badge(s, 0.9, yy, 0.62, n, c)
    txt(s, 1.75, yy+0.04, 10.7, 0.7, [[P(t, 16, INK, True, HEAD)],[P(d, 12, SLATE, False, BODY_F)]])
    yy += 0.98

# ---- 3. INTRODUCTION ----
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE); deco(s, "light")
blob(s, 9.6, 4.6, 3.4, 3.4, T_VIOLET)
chevron(s, 0.9, 1.4, 0.34, VIOLET, n=3, gap=0.3)
txt(s, 0.9, 2.0, 11.3, 0.5, [[P("INTRODUCTION", 14, TEAL_D, True, BODY_F)]])
txt(s, 0.9, 2.5, 11.5, 2.0, [
    [P("Une seule injection SQL fusionnée dans une API de paiement", 30, INK, False, HEAD)],
    [P("peut compromettre l'intégrité du système d'information de la banque.", 28, VIOLET_D, False, HEAD)],
])
rect(s, 0.92, 4.95, 1.6, 0.05, fill=GOLD)
txt(s, 0.9, 5.2, 11.0, 1.0, [
    [P("À la BTE, la revue de sécurité du code reste ", 17, BODY, False, BODY_F),
     P("entièrement manuelle", 17, GOLD, True, HEAD),
     P(".", 17, BODY, False, BODY_F)],
    [P("Aujourd'hui, aucun contrôle automatique n'arrête une faille.", 17, BODY, False, BODY_F)],
])

# ======================= CHAPITRE 1 — CONTEXTE GÉNÉRAL =======================
s = slide(); section_divider(s, 0, "Organisme d'accueil, étude de l'existant et méthodologie de travail")

# ---- Organisme & étude de l'existant ----
s = slide(); header(s, "01 · Contexte général", "Organisme d'accueil et étude de l'existant", 1)
fit_image(s, img("asis_workflow.png"), 0.7, 1.85, 6.0, 5.0, caption="Processus de revue actuel à la BTE")
txt(s, 7.1, 2.0, 5.5, 0.5, [[P("Banque de Tunisie et des Émirats", 16, INK, True, HEAD)]])
for i,(t,d) in enumerate([("Banque universelle", "particuliers, professionnels, entreprises"),
                          ("DCIO — sécurité opérationnelle", "cadre d'accueil du stage"),
                          ("Convention État tunisien · ADIA", "création en 1982, vocation élargie")]):
    rect(s, 7.1, 2.7+i*0.78, 0.16, 0.16, fill=TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, 7.4, 2.62+i*0.78, 5.2, 0.7, [[P(t+" — ", 13, INK, True, HEAD), P(d, 12, BODY, False, BODY_F)]])
card_text(s, 7.1, 5.2, 5.5, 1.3, [
    [P("Une revue intégralement manuelle", 14, RED, True, HEAD)],
    [P("aucune analyse statique ni gate CI/CD n'est greffé sur le cycle de livraison.", 12, BODY, False, BODY_F)],
], fill=PANEL, line=None)

# ---- Problématique ----
s = slide(); header(s, "01 · Contexte général", "Problématique", 2)
rect(s, 0.7, 2.0, 11.93, 4.4, fill=PANEL, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
items = [("Délai pouvant atteindre 24 h", "le temps qu'une Pull Request soit revue à la main", RED),
         ("Absence de gate CI/CD", "une Pull Request vulnérable peut être fusionnée sans blocage", RED),
         ("Qualité de revue variable", "selon la disponibilité et l'expertise du relecteur", AMBER),
         ("Angles morts", "secrets commités et dépendances vulnérables non détectés à l'œil nu", AMBER)]
for i,(t,d,c) in enumerate(items):
    col = i % 2; row = i // 2
    x = 1.1 + col*5.9; y = 2.5 + row*1.9
    icon_badge(s, x, y, 0.7, "!", c)
    txt(s, x+0.95, y+0.02, 4.7, 1.4, [[P(t, 15, INK, True, HEAD)],[P(d, 12.5, BODY, False, BODY_F)]])

# ---- Objectifs et contribution (stats animées) ----
s = slide(); header(s, "01 · Contexte général", "Objectifs et contribution", 3)
card_text(s, 0.7, 1.95, 11.93, 1.75, [
    [P("Un ", 19, BODY, False, BODY_F), P("agent IA autonome", 19, VIOLET_D, True, HEAD),
     P(" qui revoit chaque Pull Request selon l'", 19, BODY, False, BODY_F),
     P("OWASP Top 10", 19, INK, True, HEAD), P(" en ", 19, BODY, False, BODY_F),
     P("quinze minutes", 19, INK, True, HEAD), P(",", 19, BODY, False, BODY_F)],
    [P("publie les correctifs sur GitHub, et s'exécute ", 19, BODY, False, BODY_F),
     P("intégralement en local", 19, GREEN, True, HEAD),
     P(" — aucun code ne quitte la banque.", 19, BODY, False, BODY_F)],
], fill=PANEL, line=None, anchor=MSO_ANCHOR.MIDDLE)
rect(s, 0.7, 1.95, 0.13, 1.75, fill=VIOLET, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
stats = [("24 h → 15 min", "délai de revue", VIOLET_D),
         ("7 / 7", "vulnérabilités détectées", GREEN),
         ("0", "ligne de code ne sort du VPS", TEAL_D)]
cards = []
for i,(big_t,lab,c) in enumerate(stats):
    x = 0.7 + i*4.06
    sp = card_text(s, x, 4.15, 3.8, 2.0, [
        [P(big_t, 30, c, True, HEAD)], [P("", 6, BODY)], [P(lab, 13, SLATE, False, BODY_F)],
    ], anchor=MSO_ANCHOR.MIDDLE)
    for pp in sp.text_frame.paragraphs: pp.alignment = PP_ALIGN.CENTER
    rect(s, x, 4.15, 3.8, 0.12, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    cards.append(sp)
ANIM_BUILDS.append((s, cards))

# ---- Solution proposée (processus cible) ----
s = slide(); header(s, "01 · Contexte général", "Solution proposée", 4)
fit_image(s, img("tobe_workflow.png"), 0.7, 1.8, 6.6, 5.05, caption="Processus cible — BTE Security AI Agent")
txt(s, 7.7, 2.1, 5.0, 0.6, [[P("Un agent qui prend en charge tout le cycle", 16, INK, True, HEAD)]])
yy = 2.8
for t in ["Webhook signé, pipeline automatique","Cinq scanners SAST et deux LLM locaux",
          "Revue OWASP, score et verdict","Gate CI/CD (APPROVE / CHANGES / BLOCK)","Inférence locale — conformité BCT"]:
    chevron(s, 7.7, yy+0.03, 0.16, VIOLET, n=2, gap=0.13)
    txt(s, 8.1, yy-0.04, 4.6, 0.5, [[P(t, 13, BODY, False, BODY_F)]]); yy += 0.62
card_text(s, 7.7, 6.0, 5.0, 0.7, [[P("Délai : 15–25 min  vs  jusqu'à 24 h en manuel", 13, GREEN, True, HEAD)]],
          fill=PANEL, line=None, anchor=MSO_ANCHOR.MIDDLE)

# ---- Méthodologie ----
s = slide(); header(s, "01 · Contexte général", "Méthodologie de travail", 5)
fit_image(s, img("gantt_chart.png"), 0.7, 1.95, 7.7, 4.4, caption="Diagramme de Gantt — dix sprints sur cinq mois")
txt(s, 8.7, 2.05, 4.0, 0.5, [[P("Méthode Scrumban", 16, INK, True, HEAD)]])
for t,d in [("Scrum + Kanban","cadence de sprints et flux à WIP = 1"),
            ("Dix sprints / 5 mois","de la fondation à l'extension"),
            ("Système pull","chaque incident devient la carte prioritaire suivante")]:
    pass
yy=2.7
for t,d in [("Scrum + Kanban","cadence de sprints, flux à WIP = 1"),
            ("Dix sprints sur cinq mois","de la fondation Docker à l'observabilité"),
            ("Système « pull »","chaque incident de production devient la carte prioritaire")]:
    rect(s, 8.7, yy+0.04, 0.16,0.16, fill=VIOLET, shape=MSO_SHAPE.OVAL)
    txt(s, 9.0, yy-0.04, 3.8, 0.8, [[P(t, 13, INK, True, HEAD)],[P(d, 11.5, BODY, False, BODY_F)]]); yy += 1.0

# ======================= CHAPITRE 2 — ÉTAT DE L'ART =======================
s = slide(); section_divider(s, 1, "Concepts DevSecOps, spécification des besoins et choix technologiques")

# ---- DevSecOps & shift-left ----
s = slide(); header(s, "02 · État de l'art", "DevSecOps et approche shift-left", 6)
fit_image(s, img("devsecops_lifecycle.png"), 0.7, 1.9, 6.4, 4.6, caption="Cycle de vie DevSecOps")
txt(s, 7.4, 2.1, 5.2, 0.5, [[P("Intégrer la sécurité au plus tôt", 16, INK, True, HEAD)]])
for t,d in [("Security as Code","politiques de sécurité versionnées et automatisées"),
            ("Shift Left","détecter la vulnérabilité dès la Pull Request"),
            ("Le verdict devient un gate","contrôle technique conditionnant la fusion")]:
    pass
yy=2.75
for t,d in [("Security as Code","politiques de sécurité versionnées et automatisées"),
            ("Shift Left","détecter la vulnérabilité dès la Pull Request"),
            ("Le verdict comme gate","un contrôle technique conditionne la fusion")]:
    rect(s, 7.4, yy+0.04, 0.16,0.16, fill=TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, 7.7, yy-0.04, 4.9, 0.85, [[P(t, 13.5, INK, True, HEAD)],[P(d, 12, BODY, False, BODY_F)]]); yy += 1.05

# ---- Spécification des besoins ----
s = slide(); header(s, "02 · État de l'art", "Spécification des besoins", 7)
fp = card_text(s, 0.7, 2.0, 5.85, 4.4, [
    [P("Besoins fonctionnels", 16, VIOLET_D, True, HEAD)], [P("", 4, BODY)],
    [P("•  Détecter les vulnérabilités (5 scanners)", 13, BODY, False, BODY_F)],
    [P("•  Produire la revue OWASP (score, verdict, inline)", 13, BODY, False, BODY_F)],
    [P("•  Publier sur GitHub + gate CI/CD", 13, BODY, False, BODY_F)],
    [P("•  Persister chaque revue (PostgreSQL)", 13, BODY, False, BODY_F)],
    [P("•  Chat opérationnel temps réel (19 outils)", 13, BODY, False, BODY_F)],
    [P("•  Sécuriser la plateforme (HMAC, secrets exclus)", 13, BODY, False, BODY_F)],
], fill=PANEL, line=None)
fp = card_text(s, 6.78, 2.0, 5.85, 4.4, [
    [P("Besoins non fonctionnels", 16, TEAL_D, True, HEAD)], [P("", 4, BODY)],
    [P("•  Confidentialité — 100 % local (conformité BCT)", 13, BODY, False, BODY_F)],
    [P("•  Performance — 15 à 25 min par revue", 13, BODY, False, BODY_F)],
    [P("•  Résilience — circuit breaker et fallbacks", 13, BODY, False, BODY_F)],
    [P("•  Fiabilité — reprise après redémarrage du conteneur", 13, BODY, False, BODY_F)],
    [P("•  Observabilité — Prometheus / Grafana", 13, BODY, False, BODY_F)],
], fill=T_TEAL, line=None)

# ---- Choix des modèles LLM — benchmark ----
s = slide(); header(s, "02 · État de l'art", "Choix des modèles LLM — benchmark sur le matériel cible", 8)
rows = [("Modèle","Taille","Vitesse","Précision outils","Décision"),
        ("qwen2.5-coder:7b","4,7 Go","~5 tok/s","80 %","Classification + chat"),
        ("qwen2.5-coder:14b","9,0 Go","~3,2 tok/s","80 %","Analyse de sécurité"),
        ("llama3.2:3b","2,0 Go","~8 tok/s","0 %","Écarté"),
        ("granite3.1-dense:2b","1,6 Go","~8,5 tok/s","0 %","Écarté")]
tb = s.shapes.add_table(len(rows),5,Inches(0.7),Inches(1.95),Inches(11.93),Inches(3.1)).table
tb.columns[0].width=Inches(3.3); tb.columns[1].width=Inches(1.7); tb.columns[2].width=Inches(1.9)
tb.columns[3].width=Inches(2.4); tb.columns[4].width=Inches(2.63)
for r in range(len(rows)):
    for c in range(5):
        cell = tb.cell(r,c); cell.margin_left=Inches(0.12); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
        run = cell.text_frame.paragraphs[0].add_run(); run.text = rows[r][c]
        run.font.name = HEAD if (r==0 or c==0) else BODY_F
        run.font.size = Pt(12.5)
        if r==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=INK; run.font.color.rgb=WHITE; run.font.bold=True
        elif r in (1,2):
            cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0xEA,0xF7,0xF0)
            run.font.color.rgb = GREEN if c==4 else INK; run.font.bold=(c==0 or c==4)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r%2 else PANEL
            run.font.color.rgb = BODY
            if c==4: run.font.color.rgb = RED
card_text(s, 0.7, 5.35, 7.6, 1.25, [
    [P("Architecture à deux modèles : ", 13, INK, True, HEAD),
     P("le 7B classifie en ~30 s, le 14B mène l'analyse approfondie — "
       "à précision d'appel d'outils égale (80 %).", 13, BODY, False, BODY_F)],
], fill=PANEL, line=None, anchor=MSO_ANCHOR.MIDDLE)
card_text(s, 8.5, 5.35, 4.13, 1.25, [
    [P("Ollama : +22 %", 16, TEAL_D, True, HEAD)],
    [P("de débit vs LocalAI, sur GGUF identique (benchmark croisé).", 11.5, BODY, False, BODY_F)],
], fill=T_TEAL, line=None, anchor=MSO_ANCHOR.MIDDLE)

# ---- Technologies retenues (hub) ----
s = slide(); header(s, "02 · État de l'art", "Technologies retenues", 9)
triad = [("5", "Cinq scanners SAST", "Trivy · Gitleaks · Semgrep · Checkov · OSV", TEAL),
         ("IA", "Deux LLM locaux", "qwen2.5-coder 7B (classer) + 14B (analyser)", VIOLET),
         ("⟳", "Orchestration", "LangGraph — graphe d'état, reprise, routage", GOLD)]
for i,(g,t,d,c) in enumerate(triad):
    x = 0.95 + i*4.05
    rect(s, x, 2.1, 3.7, 3.7, fill=WHITE, line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    icon_badge(s, x+1.45, 2.5, 0.8, g, c)
    txt(s, x+0.3, 3.6, 3.1, 0.5, [[P(t, 16, INK, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, x+0.3, 4.2, 3.1, 1.3, [[P(d, 12.5, BODY, False, BODY_F)]], align=PP_ALIGN.CENTER)
txt(s, 0.7, 6.1, 11.93, 0.5, [[P("Sortie SAST nettoyée (−52 % de tokens) avant transmission au LLM.",
                                 13, SLATE, False, BODY_F, True)]], align=PP_ALIGN.CENTER)

# ---- Positionnement ----
s = slide(); header(s, "02 · État de l'art", "Positionnement par rapport aux solutions existantes", 9)
rows = [("Solution","Données / inférence","LLM","100 % local"),
        ("Snyk","Code envoyé au cloud","Limité","Non"),
        ("SonarQube","Auto-hébergeable","Non","Règles statiques"),
        ("Copilot Autofix","Cloud","Oui","Non"),
        ("BTE Security AI Agent","Rien ne sort du VPS","Oui","OUI")]
tb = s.shapes.add_table(len(rows),4,Inches(0.7),Inches(1.95),Inches(11.93),Inches(3.7)).table
tb.columns[0].width=Inches(3.5); tb.columns[1].width=Inches(4.4)
tb.columns[2].width=Inches(1.6); tb.columns[3].width=Inches(2.43)
for r in range(len(rows)):
    for c in range(4):
        cell = tb.cell(r,c); cell.margin_left=Inches(0.12); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
        run = cell.text_frame.paragraphs[0].add_run(); run.text = rows[r][c]
        run.font.name = HEAD if (r==0 or c==0) else BODY_F
        run.font.size = Pt(12.5 if r else 13)
        if r==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=INK; run.font.color.rgb=WHITE; run.font.bold=True
        elif r==len(rows)-1:
            cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0xEA,0xF7,0xF0)
            run.font.color.rgb = GREEN if c==3 else INK; run.font.bold=(c==0 or c==3)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r%2 else PANEL
            run.font.color.rgb = BODY
            if c==3 and rows[r][c]=="Non": run.font.color.rgb = RED
txt(s, 0.7, 5.85, 11.93, 0.5, [[P("Différenciateur : ", 14, INK, True, HEAD),
    P("l'inférence 100 % locale garantit la confidentialité bancaire et la conformité BCT.", 14, GREEN, True, HEAD)]])

# ======================= CHAPITRE 3 — CONCEPTION =======================
s = slide(); section_divider(s, 2, "Architecture de l'agent — « un cerveau, et non un simple pipeline »")

# ---- Architecture globale ----
s = slide(); header(s, "03 · Conception", "Architecture globale de la solution", 10)
fit_image(s, img("full_architecture.png"), 0.7, 1.8, 8.5, 5.0, caption="Onze conteneurs répartis en quatre couches")
layers = [("Couche entrée","nginx — webhook HMAC",TEAL),("Couche IA","FastAPI + LangGraph + Ollama",VIOLET),
          ("Couche données","PostgreSQL + Redis",INK),("Observabilité","Prometheus · Grafana",GREEN)]
for i,(t,d,c) in enumerate(layers):
    y = 2.2 + i*1.1
    rect(s, 9.5, y, 0.13, 0.85, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 9.8, y, 3.1, 0.9, [[P(t, 14, INK, True, HEAD)],[P(d, 11, BODY, False, BODY_F)]])

# ---- Manifeste 6 capacités (animées) ----
s = slide(); header(s, "03 · Conception", "Un agent autonome, et non un simple appel LLM", 11)
caps = [("1","Perception","événements webhook GitHub",TEAL),
        ("2","Raisonnement","classe la PR, route le graphe d'état",VIOLET),
        ("3","Action","commente, fixe le gate, escalade",INK),
        ("4","État et reprise","checkpoints LangGraph (PostgreSQL)",GREEN),
        ("5","Usage d'outils","dix-neuf outils réels (Docker, Prometheus…)",VIOLET),
        ("6","Auto-exploitation","scheduler, health digest, gardien disque",TEAL)]
cards=[]
for i,(n,t,d,c) in enumerate(caps):
    col=i%3; row=i//3
    x=0.7+col*4.06; y=1.95+row*2.35
    sp = card_text(s, x, y, 3.8, 2.05, [
        [P(n+"  ", 18, c, True, HEAD), P(t, 15, INK, True, HEAD)], [P("", 4, BODY)],
        [P(d, 12.5, BODY, False, BODY_F)],
    ])
    cards.append(sp)
ANIM_BUILDS.append((s, cards))

# ---- Pipeline (image + 7 étapes animées) ----
s = slide(); header(s, "03 · Conception", "Le pipeline de revue — graphe d'état LangGraph", 12)
gif_box(s, "pipeline.gif", 0.7, 1.8, 7.0, 5.0, fallback="langgraph_state_graph.png", caption="StateGraph à neuf nœuds, routage conditionnel", label="exécution du pipeline")
steps = ["webhook — HMAC + déduplication","classify — LLM 7B (~30 s)","route — scanners selon la PR",
         "scan — cinq SAST en parallèle","analyze — LLM 14B + OWASP","verdict — APPROVE/CHANGES/BLOCK",
         "publish — commentaires + gate"]
step_shapes=[]
yy=2.0
for i,st in enumerate(steps):
    sp = card_text(s, 8.0, yy, 4.7, 0.62, [
        [P(f"{i+1}  ", 13, (VIOLET_D if i in (1,4) else TEAL_D), True, HEAD),
         P(st, 12, BODY, False, BODY_F)]], anchor=MSO_ANCHOR.MIDDLE, pad=0.18)
    step_shapes.append(sp); yy += 0.68
ANIM_BUILDS.append((s, step_shapes))

# ---- Fiabilité : anti-hallucination + 2 modèles ----
s = slide(); header(s, "03 · Conception", "Fiabilité — anti-hallucination et deux modèles", 13)
fit_image(s, img("anti_hallucination_layers.png"), 0.7, 1.8, 6.0, 4.5, caption="Six couches de protection")
fit_image(s, img("two_model_architecture.png"), 7.0, 1.8, 5.6, 3.0, caption="Deux modèles : 7B pour classer, 14B pour analyser")
card_text(s, 7.0, 5.1, 5.6, 1.5, [
    [P("Garantir la fiabilité des sorties du LLM", 13, INK, True, HEAD)],
    [P("temperature=0 · num_ctx maîtrisé · garde sans-outil · prompt anti-hallucination · "
       "parser éliminant les lignes inexistantes · fusion réduisant le temps de moitié.", 12, BODY, False, BODY_F)],
], fill=PANEL, line=None)

# ======================= CHAPITRE 4 — RÉALISATION =======================
s = slide(); section_divider(s, 3, "Environnement, sécurité en action, démonstration et validation")

# ---- Environnement de travail ----
s = slide(); header(s, "04 · Réalisation", "Environnement de travail")
card_text(s, 0.7, 2.0, 5.85, 4.4, [
    [P("Environnement matériel — VPS de production", 15, INK, True, HEAD)], [P("", 5, BODY)],
    [P("•  CPU : 12 cœurs Intel Haswell (AVX2)", 13, BODY, False, BODY_F)], [P("", 3, BODY)],
    [P("•  RAM : 45 Go  ·  Disque : 290 Go", 13, BODY, False, BODY_F)], [P("", 3, BODY)],
    [P("•  GPU : aucun — inférence CPU uniquement", 13, RED, True, BODY_F)], [P("", 3, BODY)],
    [P("•  OS : Ubuntu Linux", 13, BODY, False, BODY_F)],
], fill=PANEL, line=None)
card_text(s, 6.78, 2.0, 5.85, 4.4, [
    [P("Environnement logiciel — pile open source", 15, INK, True, HEAD)], [P("", 5, BODY)],
    [P("•  Docker 29.4 — onze conteneurs", 13, BODY, False, BODY_F)], [P("", 3, BODY)],
    [P("•  Python 3.12 · FastAPI · LangGraph 1.1+", 13, BODY, False, BODY_F)], [P("", 3, BODY)],
    [P("•  PostgreSQL 16 · Redis 7 · nginx", 13, BODY, False, BODY_F)], [P("", 3, BODY)],
    [P("•  Trivy · Gitleaks · Semgrep · Checkov · OSV-Scanner", 13, BODY, False, BODY_F)],
], fill=T_TEAL, line=None)
txt(s, 0.7, 6.55, 11.93, 0.4, [[P("Une contrainte assumée : démontrer la faisabilité sur une infrastructure bancaire réelle, sans accélérateur.",
                                  12.5, SLATE, False, BODY_F, True)]], align=PP_ALIGN.CENTER)

# ---- Défense en profondeur ----
s = slide(); header(s, "04 · Réalisation", "Défense en profondeur — cinq scanners et un LLM", 14)
gif_box(s, "scanners.gif", 0.7, 1.8, 7.4, 5.0, fallback="sast_tools_overview.png", caption="Cinq outils SAST en parallèle — OWASP Top 10 2025", label="scanners en parallèle")
tools = [("Trivy","CVE et images"),("Gitleaks","secrets"),("Semgrep","OWASP code"),
         ("Checkov","IaC"),("OSV-Scanner","dépendances")]
yy=2.05
for t,d in tools:
    rect(s, 8.4, yy, 1.85, 0.42, fill=TEAL_D, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 8.4, yy+0.04, 1.85, 0.36, [[P(t, 12, WHITE, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, 10.4, yy+0.06, 2.3, 0.4, [[P(d, 12, BODY, False, BODY_F)]]); yy += 0.55
card_text(s, 8.4, 5.05, 4.25, 1.55, [
    [P("Complémentarité", 14, VIOLET_D, True, HEAD)],
    [P("Chaque scanner couvre une classe de risque — secrets, code, dépendances, IaC, images ; le modèle 14B consolide leurs résultats en une revue unique et priorisée, avec une suggestion de correction par vulnérabilité.",
       12, BODY, False, BODY_F)]], fill=T_VIOLET, line=None)

# ---- Démonstration 1/3 — le code soumis (PR #22) ----
s = slide(); header(s, "04 · Réalisation", "Démonstration — le code soumis", True)
fit_image(s, res("wire_transfer_file.png"), 0.7, 1.85, 6.5, 4.9,
          caption="wire_transfer.php — sept vulnérabilités volontaires + secrets en dur")
fit_image(s, res("requirements_file.png"), 7.55, 2.0, 5.05, 1.2, caption="requirements.txt — dépendances obsolètes")
fit_image(s, res("docker_file.png"), 7.55, 3.7, 5.05, 1.1, caption="Dockerfile — conteneur en root")
card_text(s, 7.55, 5.35, 5.05, 1.5, [
    [P("Une seule Pull Request", 13, VIOLET_D, True, HEAD)],
    [P("Trois fichiers : injection SQL, commande, SSRF, désérialisation, XSS, MD5, mot de passe en URL ; clé JWT, jeton et mot de passe en clair.",
       11, BODY, False, BODY_F)]], fill=T_VIOLET, line=None)

# ---- Démonstration 2/3 — l'agent au travail ----
s = slide(); header(s, "04 · Réalisation", "Démonstration — l'agent au travail", True)
gif_box(s, "agent_logs.gif", 0.7, 1.9, 8.4, 4.55,
        caption="Journaux de l'agent en direct — webhook → clone → classification → six scanners → analyse 14B")
txt(s, 9.15, 1.95, 3.5, 0.5, [[P("Sans intervention humaine", 14, INK, True, HEAD)]])
yy=2.65
for t in ["Webhook vérifié (HMAC-SHA256)","Classification par le modèle 7B",
          "Six scanners en parallèle","Semgrep 43 · OSV 111 · Trivy 64",
          "Analyse approfondie par le 14B","100 % local — rien ne sort du VPS"]:
    rect(s, 9.15, yy+0.05, 0.15,0.15, fill=TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, 9.4, yy-0.02, 3.25, 0.5, [[P(t, 11.5, BODY, False, BODY_F)]]); yy+=0.62

# ---- Démonstration 3/3 — verdict et revue inline (BLOCK animé) ----
s = slide(); header(s, "04 · Réalisation", "Démonstration — verdict et revue inline", True)
gif_box(s, "pr_review.gif", 0.7, 1.95, 7.3, 3.65,
        caption="Synthèse CRITIQUE, puis huit commentaires inline sur l'onglet « Files changed »")
txt(s, 8.2, 1.95, 4.45, 0.5, [[P("Huit commentaires inline", 14, INK, True, HEAD)]])
yy=2.55
for t in ["Injection SQL · commande · SSRF","Désérialisation · XSS · MD5",
          "Mot de passe transmis dans l'URL","Dockerfile sans directive USER",
          "+ une suggestion de correction par point"]:
    rect(s, 8.2, yy+0.04, 0.15,0.15, fill=TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, 8.45, yy-0.02, 4.2, 0.32, [[P(t, 11.5, BODY, False, BODY_F)]]); yy+=0.4
badge = rect(s, 8.2, 4.75, 4.45, 0.62, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
badge.shadow.inherit=False; add_shadow(badge)
tf=badge.text_frame; tf.word_wrap=True
r=tf.paragraphs[0].add_run(); r.text="VERDICT :  BLOCK"; r.font.size=Pt(18); r.font.bold=True
r.font.name=HEAD; r.font.color.rgb=WHITE; tf.paragraphs[0].alignment=PP_ALIGN.CENTER
tf.vertical_anchor=MSO_ANCHOR.MIDDLE
fit_image(s, res("pr_slack_22.png"), 8.2, 5.65, 4.45, 1.05, caption="Notification Slack — équipe prévenue")
ANIM_BUILDS.append((s, [badge]))

# ---- Bénéfices et résultats ----
s = slide(); header(s, "04 · Réalisation", "Bénéfices et résultats", 16)
rows=[("Critère","Revue manuelle","Agent IA"),
      ("Délai","jusqu'à 24 h (selon le relecteur)","15–25 min"),
      ("Couverture","variable","OWASP Top 10 systématique"),
      ("Gate CI/CD","aucun","APPROVE / CHANGES / BLOCK"),
      ("Confidentialité","—","100 % local (BCT)")]
tb=s.shapes.add_table(len(rows),3,Inches(0.7),Inches(1.95),Inches(7.6),Inches(3.7)).table
tb.columns[0].width=Inches(2.0); tb.columns[1].width=Inches(2.9); tb.columns[2].width=Inches(2.7)
for r in range(len(rows)):
    for c in range(3):
        cell=tb.cell(r,c); cell.margin_left=Inches(0.1); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
        run=cell.text_frame.paragraphs[0].add_run(); run.text=rows[r][c]
        run.font.name=HEAD if (r==0 or c==0) else BODY_F; run.font.size=Pt(12 if r else 12.5)
        if r==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=(GREEN if c==2 else INK); run.font.color.rgb=WHITE; run.font.bold=True
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb=WHITE if r%2 else PANEL; run.font.color.rgb=BODY
            if c==0: run.font.color.rgb=INK; run.font.bold=True
            if c==2: run.font.color.rgb=RGBColor(0x12,0x7A,0x46); run.font.bold=True
card_text(s, 8.6, 1.95, 4.03, 3.7, [
    [P("7 / 7", 54, GREEN, True, HEAD)],
    [P("vulnérabilités plantées", 14, INK, True, HEAD)],
    [P("détectées sur la PR #22", 12, SLATE, False, BODY_F)],
    [P("", 6, BODY)],
    [P("+ corpus de huit revues réelles", 12, BODY, False, BODY_F)],
], anchor=MSO_ANCHOR.MIDDLE)
for pp in s.shapes[-1].text_frame.paragraphs: pp.alignment=PP_ALIGN.CENTER

# ---- Observabilité 1/2 — la stack + dashboard en direct ----
s = slide(); header(s, "04 · Réalisation", "Supervision complète du VPS", 17)
gif_box(s, "grafana_live.gif", 0.7, 1.9, 7.3, 4.45, fallback="grafana_agent_dashboard.png",
        caption="Tableau de bord DevSecOps en direct — rafraîchissement 30 s", label="dashboard en direct")
txt(s, 8.35, 1.95, 4.28, 0.5, [[P("Quatre piliers open source", 14, INK, True, HEAD)]])
pillars=[("Prometheus","collecte des métriques"),("VictoriaMetrics","stockage longue durée"),
         ("Grafana","trois tableaux de bord"),("AlertManager","quinze règles d'alerte")]
yy=2.55
for t,d in pillars:
    rect(s, 8.35, yy, 1.95, 0.42, fill=TEAL_D, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 8.35, yy+0.05, 1.95, 0.34, [[P(t, 11, WHITE, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, 10.4, yy+0.07, 2.25, 0.4, [[P(d, 11, BODY, False, BODY_F)]]); yy+=0.52
card_text(s, 8.35, 4.85, 4.28, 1.55, [
    [P("28 métriques · 15 alertes · 3 dashboards", 12.5, VIOLET_D, True, HEAD)],
    [P("Métriques conteneurs via le socket Docker, après l'échec de cAdvisor sur notre version ; les deux moteurs d'inférence comparés en continu.",
       10.5, BODY, False, BODY_F)]], fill=T_VIOLET, line=None)

# ---- Observabilité 2/2 — la supervision agit (reliée à Slack) ----
s = slide(); header(s, "04 · Réalisation", "Une supervision opérationnelle, reliée à Slack", True)
fit_image(s, img("grafana_vps_host_dashboard.png"), 0.7, 1.9, 3.95, 2.1, caption="VPS — hôte")
fit_image(s, img("grafana_pr_reviews_dashboard.png"), 4.75, 1.9, 3.95, 2.1, caption="Revues de PR")
fit_image(s, img("grafana_agent_dashboard.png"), 8.75, 1.9, 3.9, 2.1, caption="Agent & LLM")
# Notification Slack : une alerte réelle déclenchée puis résolue (recadrée sur le fil de messages)
_SLK = img("slack_notification_pr_result.png")
_cl, _cr, _ct, _cb = 0.29, 0.02, 0.45, 0.18           # recadrage : retire barre latérale, haut (404) et barre de saisie
_sw, _sh = Image.open(_SLK).size
_var = (_sw*(1-_cl-_cr)) / (_sh*(1-_ct-_cb))           # ratio de la zone visible
_ph = 1.95; _pw = _ph*_var
_px, _py = 0.7, 4.6
_pcard = rect(s, _px-0.1, _py-0.1, _pw+0.2, _ph+0.2, fill=WHITE,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
round_shape(_pcard, 0.12, _pw+0.2, _ph+0.2)
_pic = s.shapes.add_picture(_SLK, Inches(_px), Inches(_py), Inches(_pw), Inches(_ph))
_pic.crop_left=_cl; _pic.crop_right=_cr; _pic.crop_top=_ct; _pic.crop_bottom=_cb
txt(s, _px, _py+_ph+0.16, _pw, 0.3,
    [[P("Canal #bte-security-ai-agent — alerte réelle, émise puis résolue", 10, SLATE, False, BODY_F, True)]],
    align=PP_ALIGN.CENTER)
_tx = _px + _pw + 0.55
txt(s, _tx, 4.7, 12.63-_tx, 0.55, [[P("La supervision agit — elle ne décore pas", 16, INK, True, HEAD)]])
for i, t in enumerate([
        "Une alerte réelle se déclenche : CPU de l'hôte à 99 %, notifiée sur Slack",
        "L'agent confirme ensuite la résolution automatique de l'incident",
        "Le même canal reçoit le verdict de chaque revue de Pull Request"]):
    rect(s, _tx, 5.4+i*0.52+0.04, 0.15, 0.15, fill=TEAL_D, shape=MSO_SHAPE.OVAL)
    txt(s, _tx+0.27, 5.37+i*0.52, 12.63-(_tx+0.27), 0.5, [[P(t, 12, BODY, False, BODY_F)]])

# ---- Assistant conversationnel & autonomie ----
s = slide(); header(s, "04 · Réalisation", "Assistant conversationnel et opérations autonomes", 18)
gif_box(s, "chat_demo.gif", 0.7, 1.85, 6.4, 4.9, fallback="inter_react.png", caption="Boucle ReAct — dix-neuf outils de supervision", label="chat répondant en direct")
txt(s, 7.4, 2.05, 5.2, 0.6, [[P("« Donne-moi la santé de mon VPS »", 16, VIOLET_D, True, HEAD)]])
txt(s, 7.4, 2.65, 5.2, 0.9, [[P("L'agent invoque un outil réel ", 13, BODY, False, BODY_F),
    P("(vps_status)", 12, INK, True, MONO), P(" et répond avec la valeur mesurée — jamais inventée.", 13, BODY, False, BODY_F)]])
txt(s, 7.4, 3.75, 5.2, 0.5, [[P("Et il s'auto-exploite :", 14, INK, True, HEAD)]])
for i,t in enumerate(["Scheduler autonome (tâches asyncio)","Health digest Slack quotidien (09:00 UTC)",
                      "Gardien disque auto-réparateur à 90 %"]):
    rect(s, 7.4, 4.3+i*0.62+0.03, 0.16,0.16, fill=GOLD, shape=MSO_SHAPE.OVAL)
    txt(s, 7.7, 4.3+i*0.62-0.03, 4.9, 0.5, [[P(t, 13, BODY, False, BODY_F)]])

# ---- Tests et validation ----
s = slide(); header(s, "04 · Réalisation", "Tests et validation de la plateforme")
tests = [("HMAC-SHA256", "signature invalide rejetée en 403 ; signature valide acceptée en 202", TEAL),
         ("Déduplication Redis", "deux webhooks identiques ne déclenchent qu'un seul pipeline (SET NX)", VIOLET),
         ("Circuit breaker LLM", "Ollama indisponible : repli sur la classification par extensions, sans blocage", GOLD),
         ("Gardien disque", "au-delà de 90 % : nettoyage Docker automatique et notification Slack", INK)]
for i,(t,d,c) in enumerate(tests):
    col = i % 2; row = i // 2
    x = 0.7 + col*6.08; y = 1.95 + row*1.7
    sp = card_text(s, x, y, 5.85, 1.5, [
        [P("✓  ", 15, GREEN, True, HEAD), P(t, 14.5, INK, True, HEAD)],
        [P(d, 12, BODY, False, BODY_F)],
    ])
    rect(s, x, y, 0.1, 1.5, fill=c)
badge = rect(s, 0.7, 5.6, 11.93, 0.85, fill=RGBColor(0xEA,0xF7,0xF0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 0.7, 5.78, 11.93, 0.5, [[P("Checklist de validation finale : ", 14, INK, True, HEAD),
    P("quatorze vérifications de bout en bout, toutes validées avant la livraison.", 14, GREEN, True, HEAD)]],
    align=PP_ALIGN.CENTER)

# ======================= CHAPITRE 5 — CONCLUSION =======================
s = slide(); section_divider(s, 4, "Limites, perspectives et bilan du projet")

# ---- Limites et perspectives ----
s = slide(); header(s, "05 · Conclusion", "Limites assumées et perspectives", 20)
card_text(s, 0.7, 2.0, 5.85, 4.3, [
    [P("Limites assumées", 17, AMBER, True, HEAD)], [P("", 6, BODY)],
    [P("•  Inférence CPU : 15–25 min/PR — le prix de la confidentialité", 13, BODY, False, BODY_F)], [P("", 4, BODY)],
    [P("•  Corpus d'évaluation restreint (jeu contrôlé + 8 revues réelles)", 13, BODY, False, BODY_F)], [P("", 4, BODY)],
    [P("•  Diff = entrée non fiable → findings scanners déterministes,", 13, BODY, False, BODY_F)],
    [P("    ancrage renforcé du verdict prévu (défense en profondeur)", 13, BODY, False, BODY_F)],
], fill=PANEL, line=None)
card_text(s, 6.78, 2.0, 5.85, 4.3, [
    [P("Perspectives", 17, GREEN, True, HEAD)], [P("", 6, BODY)],
    [P("•  Migration GPU : 30–50 tok/s → 2–5 min/PR (~3000 €)", 13, BODY, False, BODY_F)], [P("", 4, BODY)],
    [P("•  Passage à l'échelle de tous les dépôts BTE", 13, BODY, False, BODY_F)], [P("", 4, BODY)],
    [P("•  Tableau de bord de risque pour le RSSI", 13, BODY, False, BODY_F)],
], fill=T_TEAL, line=None)

# ---- Conclusion & remerciements ----
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE); deco(s, "divider")
chevron(s, 0.9, 1.2, 0.3, VIOLET, n=3, gap=0.27)
txt(s, 0.9, 1.7, 11.5, 0.4, [[P("CONCLUSION", 14, TEAL_D, True, BODY_F)]])
txt(s, 0.9, 2.15, 11.6, 1.8, [
    [P("Un agent IA agentique, déployé en production à la BTE,", 25, INK, False, HEAD)],
    [P("qui sécurise chaque Pull Request ", 25, INK, False, HEAD),
     P("sans qu'aucun code ne quitte la banque.", 25, VIOLET_D, False, HEAD)],
])
for i,(b,l,c) in enumerate([("24 h → 15 min","délai",VIOLET_D),("7 / 7","détection",GREEN),("100 % local","conformité BCT",TEAL_D)]):
    x=0.9+i*4.0
    card_text(s, x, 4.2, 3.7, 1.3, [[P(b,22,c,True,HEAD)],[P(l,12,SLATE,False,BODY_F)]], anchor=MSO_ANCHOR.MIDDLE)
    for pp in s.shapes[-1].text_frame.paragraphs: pp.alignment=PP_ALIGN.CENTER
txt(s, 0.9, 5.9, 11.5, 0.6, [[P("Plateforme remise à la BTE   ·   Merci de votre attention", 17, GOLD, True, HEAD)]])

# ===================== animations =====================
for sl in prs.slides:
    add_fade_transition(sl)
for sl, shapes in ANIM_BUILDS:
    add_fade_in_build(sl, shapes)

# ===================== discours (notes orateur) =====================
SPEECH = [
# 1 — Titre (≈ 30 s)
"(≈ 30 s) Monsieur le Président du jury, Madame, Messieurs les membres du jury, bonjour. "
"Je m'appelle Ghaith Ferchichi et j'ai l'honneur de vous présenter mon projet de fin d'études, "
"réalisé au sein de la Banque de Tunisie et des Émirats, sous l'encadrement de Monsieur Kamel "
"Kaouech, côté banque, et de Madame Ghayet El Mouna Zhioua, côté ISI. Il s'intitule BTE Security "
"AI Agent : un agent d'intelligence artificielle pour la revue automatisée de code.",
# 2 — Sommaire (≈ 10 s)
"(≈ 10 s) Mon exposé suit la structure du rapport : le contexte général, l'état de l'art et les "
"choix technologiques, la conception de l'agent, sa réalisation et ses résultats, puis la "
"conclusion et les perspectives.",
# 3 — Introduction (≈ 35 s)
"(≈ 35 s) Partons de trois scénarios concrets : une injection SQL fusionnée dans une API de "
"paiement, un secret d'authentification commité par erreur dans un dépôt, une image Docker construite "
"sur une bibliothèque touchée par une faille critique. Chacune de ces failles peut compromettre le "
"système d'information de la banque. Et aujourd'hui, aucun contrôle automatique ne les arrête : "
"la revue de sécurité du code reste entièrement manuelle. C'est ce constat qui motive ce travail.",
# 4 — Divider 1 (≈ 5 s)
"(≈ 5 s) Commençons par le contexte général.",
# 5 — Organisme & étude de l'existant (≈ 18 s)
"(≈ 18 s) Mon stage s'est déroulé à la Banque de Tunisie et des Émirats, à la Direction Centrale "
"de l'Informatique et de l'Organisation, au sein de l'équipe sécurité opérationnelle. Le "
"processus existant, à gauche : à chaque Pull Request, un relecteur examine le diff à la main, "
"commente, puis fusionne. Aucun scanner, aucun gate de pipeline — tout repose sur l'humain.",
# 6 — Problématique (≈ 27 s)
"(≈ 27 s) Quatre limites en découlent. Le délai : jusqu'à vingt-quatre heures. L'absence de "
"blocage technique : une Pull Request critique peut être fusionnée sans obstacle. Une qualité de "
"revue inégale, qui dépend de la disponibilité et de l'expertise du relecteur. Et deux angles "
"morts : les secrets commités et les dépendances vulnérables, pour lesquels l'œil humain n'est "
"pas le bon outil.",
# 7 — Objectifs et contribution (builds : 3 clics) (≈ 40 s)
"(≈ 40 s) Notre contribution tient en une phrase : un agent d'intelligence artificielle autonome "
"qui revoit chaque Pull Request selon la grille OWASP Top 10, publie ses correctifs directement "
"sur GitHub, et s'exécute intégralement en local — aucun code ne quitte la banque. Trois chiffres "
"résument le résultat. [clic] Le délai de revue passe de vingt-quatre heures à une quinzaine de "
"minutes. [clic] Sur la validation contrôlée, sept vulnérabilités sur sept sont détectées, sans "
"faux positif dans le fichier modifié. [clic] Et cent pour cent de l'inférence reste sur le VPS "
"de la banque, conformément aux exigences de la Banque Centrale de Tunisie.",
# 8 — Solution proposée (≈ 30 s)
"(≈ 30 s) Voici le processus cible. GitHub émet un webhook signé HMAC-SHA256 ; l'agent vérifie la "
"signature, élimine les doublons, classifie la Pull Request en une trentaine de secondes, exécute "
"en parallèle les scanners pertinents, puis consolide leurs résultats avec un modèle de langage "
"local. Il publie sur GitHub la revue complète : score de risque, commentaires ligne par ligne et "
"verdict — APPROVE, REQUEST_CHANGES ou BLOCK — posé comme statut de commit. C'est ce statut qui "
"conditionne la fusion, laquelle reste déclenchée par un humain.",
# 9 — Méthodologie (≈ 20 s)
"(≈ 20 s) Le projet a suivi la méthode Scrumban, retenue après comparaison avec Scrum, Kanban et "
"XP : la cadence et les jalons de Scrum, avec la discipline de flux de Kanban. Dix sprints de "
"deux semaines sur cinq mois, et une étiquette Expedite pour traiter un incident de production "
"sans replanifier le sprint en cours.",
# 10 — Divider 2 (≈ 5 s)
"(≈ 5 s) Passons à l'état de l'art et aux choix technologiques.",
# 11 — DevSecOps (≈ 18 s)
"(≈ 18 s) Le cadre conceptuel est le DevSecOps et le shift-left : détecter les vulnérabilités au "
"plus tôt, dès la Pull Request, là où une faille coûte presque rien à corriger. Notre apport "
"matérialise ce principe : le verdict de l'agent devient un gate technique qui conditionne la "
"fusion.",
# 12 — Besoins (≈ 18 s)
"(≈ 18 s) Six besoins fonctionnels : détecter les vulnérabilités, produire la revue OWASP avec "
"score et verdict, publier sur GitHub avec le gate CI/CD, persister chaque revue, offrir un chat "
"opérationnel, et sécuriser la plateforme elle-même. Côté non fonctionnel, cinq contraintes : la "
"confidentialité — l'inférence reste locale —, la performance, la résilience, la fiabilité et "
"l'observabilité.",
# 13 — Benchmark (≈ 32 s)
"(≈ 32 s) Le choix des modèles repose non sur la littérature, mais sur un banc d'essai sur le "
"matériel cible, prompt système complet. Deux candidats éliminés : llama 3.2 trois milliards, "
"saturé par le seul prompt système, et granite, dont l'appel d'outils est incompatible. Restent "
"qwen2.5-coder 7B et 14B, à quatre-vingts pour cent de précision. D'où deux modèles : le 7B "
"classifie, le 14B analyse. Le banc a aussi confirmé Ollama : vingt-deux pour cent de débit de "
"plus que LocalAI, à poids identiques.",
# 14 — Technologies retenues (≈ 30 s)
"(≈ 30 s) Trois briques portent l'ensemble. Cinq scanners spécialisés : Trivy pour les CVE, "
"Gitleaks pour les secrets, Semgrep pour le code applicatif, Checkov pour l'infrastructure, "
"OSV-Scanner pour les dépendances. Deux modèles servis localement par Ollama. Et l'orchestrateur "
"LangGraph, préféré à Prefect et Celery pour sa persistance d'état native dans PostgreSQL : un "
"pipeline de vingt minutes survit au redémarrage de son conteneur. Détail qui compte : les "
"sorties des scanners sont nettoyées avant transmission au modèle — cinquante-deux pour cent de "
"tokens en moins — pour préserver la fenêtre de contexte.",
# 15 — Positionnement (≈ 22 s)
"(≈ 22 s) Pourquoi pas un outil du marché ? Aucun ne combine raisonnement par modèle de langage "
"et inférence entièrement locale. Snyk et Copilot envoient le code à un tiers — "
"inacceptable en banque. SonarQube est auto-hébergeable mais purement statique, sans explication. "
"Notre contribution unit ces deux mondes, sans qu'une ligne ne sorte du VPS.",
# 16 — Divider 3 (≈ 5 s)
"(≈ 5 s) Venons-en à la conception.",
# 17 — Architecture (≈ 25 s)
"(≈ 25 s) La plateforme compte onze conteneurs Docker en quatre couches. L'entrée : nginx, seul "
"point exposé, qui authentifie et transmet les webhooks. La couche IA : l'agent FastAPI, qui "
"orchestre le graphe LangGraph et les scanners, et Ollama, accessible uniquement depuis le réseau "
"Docker interne. Les données : PostgreSQL pour la base de connaissances et les checkpoints, Redis "
"pour la déduplication et le cache. Et l'observabilité, sur laquelle je reviendrai.",
# 18 — Manifeste (builds : 6 clics) (≈ 35 s)
"(≈ 35 s) Ce point fonde le titre du projet : il s'agit d'un agent, et non d'un simple appel à un "
"modèle. Six capacités le démontrent. [clics] Il perçoit son environnement — webhooks, métriques, "
"alertes. Il raisonne : il classifie chaque Pull Request et route son graphe d'état en "
"conséquence. Il agit : il commente, il bloque, il escalade vers Slack. Il persiste : son état "
"est sauvegardé après chaque nœud, et il reprend exactement où il s'était arrêté. Il manipule "
"dix-neuf outils réels. Et il s'auto-exploite : gardien disque et bilan de santé quotidien.",
# 19 — Pipeline (builds : 7 clics) (≈ 42 s)
"(≈ 42 s) Le cœur du système est un graphe d'état à neuf nœuds. Suivons une Pull Request. [clics] "
"Le webhook est vérifié, les doublons éliminés ; le dépôt est cloné, le diff généré localement. Le 7B "
"classifie en cinq catégories, et le routage adapte l'analyse : une PR de documentation saute les "
"scanners, un Dockerfile déclenche le scan d'image. Les scanners tournent en parallèle, chacun "
"isolé — si l'un échoue, le pipeline continue. Le 14B mène l'analyse OWASP, rend le verdict, "
"escalade sur Slack si besoin, puis publie. Optimisation clé : fusionner les deux appels du 14B a "
"fait passer le pipeline de trente-cinquante minutes à quinze-vingt-cinq.",
# 20 — Fiabilité (≈ 35 s)
"(≈ 35 s) Reste la question légitime : peut-on faire confiance à un modèle de langage ? Six "
"couches de protection répondent. Trois agissent sur le modèle : température proche de zéro, "
"fenêtre de contexte dimensionnée, génération plafonnée. Trois relèvent du contrôle applicatif : "
"une garde qui force l'appel d'outil pour toute donnée en temps réel, des règles "
"anti-hallucination dans le prompt, et l'interdiction de citer une valeur absente des "
"observations. Pour la revue de code s'y ajoute un parser dédié, qui écarte tout commentaire "
"visant une ligne inexistante du diff. Résultat : aucune erreur de ligne publiée sur GitHub.",
# 21 — Divider 4 (≈ 5 s)
"(≈ 5 s) Passons à la réalisation et aux résultats.",
# 22 — Environnement (≈ 20 s)
"(≈ 20 s) Le terrain d'exécution est volontairement contraint : un VPS de production de douze "
"cœurs Haswell, quarante-cinq giga-octets de mémoire, sans GPU — c'est le prix de la "
"confidentialité, et nous l'assumons. Cette contrainte a imposé un vrai tuning : flash attention "
"et quantisation du cache KV pour tenir le modèle 14B en mémoire. La pile est intégralement open "
"source.",
# 23 — Défense en profondeur (≈ 30 s)
"(≈ 30 s) La détection suit une logique de défense en profondeur : chaque scanner couvre une "
"classe de risque — secrets, code, dépendances, infrastructure, images — et le modèle 14B "
"consolide l'ensemble en une revue unique, priorisée, avec une suggestion de correction par "
"point. Les findings des scanners restent déterministes : ils sont persistés en base quel que "
"soit le texte produit par le modèle.",
# 24a — Démonstration 1/3 : le code soumis (≈ 25 s)
"(≈ 25 s) La démonstration part d'une seule Pull Request, dans un dépôt de test. Trois fichiers : "
"un module de virement en PHP qui concentre sept vulnérabilités types — injection SQL, injection "
"de commande, SSRF, désérialisation non sécurisée, XSS, hachage MD5 et mot de passe transmis dans "
"l'URL — avec en prime trois secrets codés en dur. À côté, un requirements aux dépendances "
"obsolètes et un Dockerfile qui tourne en root. Le terrain de jeu idéal pour l'agent.",
# 24b — Démonstration 2/3 : l'agent au travail (≈ 30 s)
"(≈ 30 s) Le développeur crée sa branche, pousse, ouvre la Pull Request — et n'a plus rien à "
"faire. À droite, les journaux de l'agent défilent en direct : le webhook signé est vérifié, le "
"modèle 7B classe la Pull Request, puis six scanners s'exécutent en parallèle — Semgrep remonte "
"quarante-trois alertes, OSV cent onze vulnérabilités de dépendances, Trivy soixante-quatre. Tout "
"se passe en local, sans qu'aucune ligne ne quitte la banque. Le modèle 14B prend alors le relais "
"pour l'analyse approfondie.",
# 24c — Démonstration 3/3 : verdict et revue inline (≈ 26 s)
"(≈ 26 s) Le résultat est publié directement sur la Pull Request. Une synthèse en tête — risque "
"CRITIQUE, verdict BLOCK — puis huit commentaires inline, chacun ancré sur la ligne exacte du "
"diff et porteur d'une suggestion de correction. [clic] Le statut du commit passe au rouge : la "
"fusion est techniquement bloquée, et l'équipe est prévenue sur Slack. L'agent n'a pas seulement "
"commenté — il a arrêté le code vulnérable à la porte.",
# 25 — Bénéfices & résultats (≈ 25 s)
"(≈ 25 s) Le tableau résume le passage du manuel à l'agent : un délai divisé par près de cent, "
"une couverture OWASP systématique, un gate effectif là où rien ne bloquait, une traçabilité "
"complète en base. Sur le jeu contrôlé, sept sur sept, sans faux positif dans le fichier modifié. "
"Au-delà de la démonstration, huit revues réelles sont persistées en base, et leurs verdicts "
"suivent le risque.",
# 26a — Observabilité 1/2 : la stack (≈ 18 s)
"(≈ 18 s) L'agent voit l'intégralité du VPS : vingt-huit métriques, quinze règles d'alerte et "
"trois tableaux de bord Grafana, rafraîchis toutes les trente secondes. Faute de compatibilité de "
"cAdvisor avec notre version de Docker, les métriques de conteneurs passent directement par le "
"socket Docker, et les deux moteurs d'inférence sont comparés en continu.",
# 26b — Observabilité 2/2 : la supervision agit (≈ 17 s)
"(≈ 17 s) Et cette supervision n'est pas décorative : elle agit. Quand la charge processeur de "
"l'hôte grimpe, une alerte se déclenche et part sur Slack ; une fois la situation rétablie, "
"l'agent confirme lui-même la résolution. Le même canal reçoit le verdict de chaque revue. La "
"boucle observation, alerte, action est ainsi bouclée, en conditions réelles.",
# 27 — Chat & autonomie (≈ 25 s)
"(≈ 25 s) Un assistant conversationnel doté de dix-neuf outils permet d'interroger "
"l'infrastructure en langage naturel, jusqu'à la base de connaissances sécurité, sans écrire de "
"SQL. Chaque question sur une donnée réelle déclenche l'appel d'outil correspondant : la réponse "
"cite la valeur mesurée, jamais inventée. Et l'agent s'auto-exploite : gardien disque toutes les "
"trente minutes, nettoyage automatique à quatre-vingt-dix pour cent, bilan de santé publié chaque "
"matin sur Slack.",
# 28 — Tests et validation (≈ 20 s)
"(≈ 20 s) La validation finale a couvert les mécanismes critiques : signature invalide rejetée "
"avant tout traitement, webhook dupliqué ignoré, circuit breaker avec classification de repli si "
"le modèle est indisponible, gardien disque vérifié en conditions réelles. Au total, quatorze "
"vérifications de bout en bout, toutes validées avant la livraison.",
# 29 — Divider 5 (≈ 5 s)
"(≈ 5 s) J'en viens à la conclusion.",
# 30 — Limites assumées & perspectives (≈ 34 s)
"(≈ 34 s) Deux limites assumées et un point de vigilance. L'inférence sur CPU borne la revue à "
"quinze-vingt-cinq minutes : c'est le prix de la confidentialité. Le corpus d'évaluation reste "
"restreint — un jeu contrôlé et huit revues réelles. Et le diff est par nature une entrée non "
"fiable : l'architecture y oppose déjà les findings déterministes des scanners, qu'aucun texte ne "
"peut effacer. Trois perspectives : la migration GPU — environ trois mille euros — ramènerait "
"chaque revue à deux à cinq minutes, sans toucher à l'architecture ni à la confidentialité ; le "
"passage à l'échelle de tous les dépôts de la banque ; et un tableau de bord de risque pour le "
"RSSI.",
# 31 — Conclusion (≈ 30 s)
"(≈ 30 s) En conclusion : parti d'un VPS vierge, ce projet livre un agent d'intelligence "
"artificielle complet, déployé en production à la BTE, qui sécurise chaque Pull Request sans "
"qu'aucun code ne quitte la banque. La plateforme est remise à la banque et constitue un socle "
"concret pour sa transformation DevSecOps. Je vous remercie de votre attention et me tiens à "
"votre disposition pour répondre à vos questions.",
]
_slides = list(prs.slides)
if len(SPEECH) != len(_slides):
    print(f"ATTENTION : {len(SPEECH)} notes pour {len(_slides)} slides")
for sl, note_text in zip(_slides, SPEECH):
    sl.notes_slide.notes_text_frame.text = note_text

prs.save(OUT)
print("OK ->", OUT, "| slides:", len(prs.slides._sldIdLst), "| builds:", len(ANIM_BUILDS), "| notes:", len(SPEECH))
