#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génère la présentation de soutenance — BTE Security AI Agent.
Thème clair "Secure Intelligence (Light)". Réf : pptx-plan.md
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.shapes.freeform import FreeformBuilder
from PIL import Image

IMG = "/opt/devsecops/MASTER_PFE/img"
OUT = "/opt/devsecops/BTE_Security_AI_Agent_Soutenance_Academic.pptx"

# ---------- Palette : Academic (PhD) — paper + ink + teal ----------
# Noms de variables identiques au thème tech : chaque bloc de slide se
# restyle automatiquement ; seules les valeurs RGB changent.
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)   # cartes (blanc sur crème)
CREAM   = RGBColor(0xFA, 0xF7, 0xF1)   # canvas
MIST    = RGBColor(0xF3, 0xEE, 0xE3)   # panneau chaud
NAVY    = RGBColor(0x1E, 0x2A, 0x44)   # encre / titres
SLATE_T = RGBColor(0x33, 0x40, 0x5C)   # corps de texte
CYAN_D  = RGBColor(0x15, 0x7A, 0x82)   # teal — traits/labels (texte OK)
CYAN    = RGBColor(0x2B, 0xA6, 0xAE)   # teal — remplissages
INDIGO  = RGBColor(0x3E, 0x4C, 0x7A)   # bleu ardoise profond
VIOLET  = RGBColor(0x6B, 0x5B, 0xA6)   # prune sobre
GOLD    = RGBColor(0xC2, 0x91, 0x2E)   # ocre — identité BTE
SLATE   = RGBColor(0x6B, 0x62, 0x56)   # gris chaud secondaire
HAIR    = RGBColor(0xE3, 0xDC, 0xCD)   # filet chaud
GREEN   = RGBColor(0x2F, 0x7D, 0x52)   # APPROVE (sobre)
AMBER   = RGBColor(0xC5, 0x7A, 0x12)   # REQUEST_CHANGES (sobre)
RED     = RGBColor(0xB2, 0x3A, 0x48)   # BLOCK (brique sobre)

HEAD = "Georgia"          # titres serif (académique)
BODY = "Segoe UI"         # corps sans
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------
def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        add_shadow(sp)
    return sp

def add_shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = spPr.makeelement(qn('a:outerShdw'),
        {'blurRad':'90000','dist':'38100','dir':'5400000','rotWithShape':'0'})
    clr = spPr.makeelement(qn('a:srgbClr'), {'val':'0B2545'})
    alpha = spPr.makeelement(qn('a:alpha'), {'val':'22000'})
    clr.append(alpha); sh.append(clr); el.append(sh); spPr.append(el)

def set_gradient(sp, c1, c2, angle_deg=30):
    """Linear gradient fill c1->c2 (RGBColor)."""
    spPr = sp._element.spPr
    for t in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        e = spPr.find(qn(t))
        if e is not None: spPr.remove(e)
    ang = int(angle_deg*60000)
    xml = (
      '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
      '<a:gsLst>'
      f'<a:gs pos="0"><a:srgbClr val="{str(c1)}"/></a:gs>'
      f'<a:gs pos="100000"><a:srgbClr val="{str(c2)}"/></a:gs>'
      '</a:gsLst>'
      f'<a:lin ang="{ang}" scaled="1"/>'
      '</a:gradFill>'
    )
    from pptx.oxml import parse_xml
    grad = parse_xml(xml)
    ln = spPr.find(qn('a:ln'))
    if ln is not None: ln.addprevious(grad)
    else: spPr.append(grad)

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, sp_after=4):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold, font, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        for (t, sz, col, bold, fnt, *rest) in para:
            ital = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = ital
            r.font.name = fnt; r.font.color.rgb = col
    return tb

def P(t, sz, col=SLATE_T, bold=False, fnt=BODY, ital=False):
    return (t, sz, col, bold, fnt, ital)

def header(s, kicker, title, num=None):
    """En-tête académique : kicker teal + titre serif + filet fin + carré ocre."""
    rect(s, 0, 0, 13.333, 7.5, fill=CREAM)            # canvas papier
    rect(s, 0, 0.5, 0.012, 0.95, fill=GOLD)           # mince trait vertical ocre
    if kicker:
        txt(s, 0.62, 0.46, 11.0, 0.35, [[P(kicker.upper(), 12, CYAN_D, True, BODY)]])
    txt(s, 0.62, 0.74, 12.1, 0.9, [[P(title, 27, NAVY, True, HEAD)]])
    rect(s, 0.64, 1.55, 0.16, 0.16, fill=GOLD)        # petit carré géométrique
    rect(s, 0.86, 1.62, 1.5, 0.02, fill=CYAN_D)       # filet fin teal
    footer(s, num)

def footer(s, num=None):
    rect(s, 0.62, 7.06, 12.1, 0.012, fill=HAIR)
    txt(s, 0.62, 7.12, 8.0, 0.3,
        [[P("BTE Security AI Agent", 9, SLATE, False, BODY),
          P("  ·  Ghaith Ferchichi", 9, SLATE, False, BODY)]])
    if num is not None:
        rect(s, 12.0, 7.17, 0.10, 0.10, fill=CYAN_D)   # petit carré teal
        txt(s, 12.15, 7.10, 0.5, 0.3, [[P(f"{num:02d}", 10, NAVY, True, HEAD)]],
            align=PP_ALIGN.LEFT)

def fit_image(s, path, bx, by, bw, bh, card=True, caption=None, align="center"):
    """Place image preserving aspect ratio inside box, in a white card."""
    iw, ih = Image.open(path).size
    ar = iw/ih; bar = bw/bh
    if ar > bar:
        w = bw; h = bw/ar
    else:
        h = bh; w = bh*ar
    x = bx + (bw-w)/2; y = by + (bh-h)/2
    if card:
        pad = 0.10
        c = rect(s, x-pad, y-pad, w+2*pad, h+2*pad, fill=WHITE,
                 line=HAIR, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if caption:
        txt(s, bx, by+bh+0.06, bw, 0.3, [[P(caption, 10, SLATE, False, BODY, True)]],
            align=PP_ALIGN.CENTER)
    return (x, y, w, h)

def chip(s, x, y, w, label, color, txtcol=WHITE, h=0.34):
    c = rect(s, x, y, w, h, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, y+0.02, w, h-0.02, [[P(label, 11, txtcol, True, HEAD)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c

def img(name):
    for sub in ("chapter_1","chapter_2","chapter_3","chapter_4","template"):
        p = os.path.join(IMG, sub, name)
        if os.path.exists(p): return p
    raise FileNotFoundError(name)

# ================= SLIDES =================

# ---- 1. TITLE ----
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=CREAM)
rect(s, 0.9, 1.95, 11.5, 0.02, fill=NAVY)      # filet fin haut
rect(s, 0.9, 1.95, 2.4, 0.02, fill=GOLD)       # segment ocre superposé
# logos — BTE (entreprise) à gauche ; ISI + UTM (université) à droite
def logo_w(nm, H):
    iw, ih = Image.open(img(nm)).size
    return H*(iw/ih)
try:
    s.shapes.add_picture(img("logo_BTE.png"), Inches(0.7), Inches(0.55), height=Inches(0.60))
except Exception:
    pass
Hu = 0.66
try:
    wu = logo_w("Logo_UTM.png", Hu)
    xu = 12.63 - wu
    s.shapes.add_picture(img("Logo_UTM.png"), Inches(xu), Inches(0.5), height=Inches(Hu))
    wi = logo_w("LogoISI.png", 0.58)
    s.shapes.add_picture(img("LogoISI.png"), Inches(xu - 0.5 - wi), Inches(0.56), height=Inches(0.58))
except Exception:
    pass
txt(s, 0.9, 2.3, 11.5, 0.4, [[P("RAPPORT DE STAGE DE FIN D'ÉTUDES  ·  MASTÈRE PROFESSIONNEL SSI", 13, CYAN_D, True, BODY)]])
txt(s, 0.9, 2.8, 11.6, 1.6, [
    [P("BTE Security AI Agent", 44, NAVY, True, HEAD)],
    [P("Agent IA agentique pour la revue automatisée de sécurité du code", 21, SLATE_T, False, BODY, True)],
])
rect(s, 0.92, 4.5, 0.16, 0.16, fill=GOLD)          # petit carré
rect(s, 1.14, 4.57, 2.0, 0.02, fill=CYAN_D)        # filet fin teal
txt(s, 0.9, 4.8, 11.5, 1.2, [
    [P("Réalisé par  ", 14, SLATE, False, BODY), P("Ghaith FERCHICHI", 14, NAVY, True, HEAD)],
    [P("Encadrant pro : ", 12, SLATE, False, BODY), P("M. Kamel KAOUECH", 12, SLATE_T, False, BODY),
     P("     Encadrante ISI : ", 12, SLATE, False, BODY), P("Mme Ghayet El Mouna ZHIOUA", 12, SLATE_T, False, BODY)],
    [P("Banque de Tunisie et des Émirats  ·  Année universitaire 2025–2026", 12, SLATE, False, BODY, True)],
])
rect(s, 0.9, 6.7, 11.5, 0.02, fill=NAVY)       # filet fin bas
rect(s, 9.99, 6.7, 2.41, 0.02, fill=GOLD)      # segment ocre superposé

# ---- 2. SOMMAIRE ----
s = slide(); header(s, "Plan", "Sommaire", 1)
items = [
    ("01", "Contexte & problématique", CYAN_D),
    ("02", "Objectifs & positionnement", CYAN_D),
    ("03", "Architecture de la solution", NAVY),
    ("04", "L'IA agentique au cœur du système", INDIGO),
    ("05", "Pipeline de sécurité & bénéfices", GREEN),
    ("06", "Démonstration & résultats", NAVY),
    ("07", "Limites & perspectives", SLATE),
]
y = 2.0
for i,(n,t,c) in enumerate(items):
    col = 0 if i < 4 else 1
    yy = 2.0 + (i % 4)*1.15
    xx = 0.8 + col*6.2
    rect(s, xx, yy, 0.9, 0.9, fill=MIST, line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, xx, yy+0.22, 0.9, 0.5, [[P(n, 22, c, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, xx+1.1, yy+0.22, 5.0, 0.6, [[P(t, 15, NAVY, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)

# ---- 3. ACCROCHE (hook) ----
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=CREAM)
rect(s, 1.0, 1.55, 0.06, 3.4, fill=GOLD)          # filet vertical ocre (citation)
txt(s, 1.4, 1.55, 11.0, 0.5, [[P("LE PROBLÈME", 14, CYAN_D, True, BODY)]])
txt(s, 1.4, 2.2, 11.0, 3.0, [
    [P("Une seule injection SQL fusionnée dans une API de paiement", 30, NAVY, True, HEAD)],
    [P("peut compromettre toute la banque.", 30, CYAN_D, True, HEAD)],
])
txt(s, 1.4, 4.7, 11.0, 1.2, [
    [P("Aujourd'hui à la BTE, la revue de sécurité est ", 18, SLATE_T, False, BODY),
     P("manuelle", 18, GOLD, True, HEAD),
     P(", prend ", 18, SLATE_T, False, BODY),
     P("jusqu’à 24 h", 18, GOLD, True, HEAD),
     P(",", 18, SLATE_T, False, BODY)],
    [P("et dépend entièrement du relecteur de garde.", 18, SLATE_T, False, BODY)],
])
rect(s, 1.0, 6.55, 11.5, 0.02, fill=HAIR)

# ---- 4. CONTEXTE AS-IS ----
s = slide(); header(s, "01 · Contexte", "Le processus de revue actuel (AS-IS)", 2)
fit_image(s, img("asis_workflow.png"), 0.7, 1.8, 6.2, 5.0, caption="Processus AS-IS à la BTE")
txt(s, 7.4, 2.0, 5.4, 0.5, [[P("Limites identifiées", 17, NAVY, True, HEAD)]])
lims = [
    ("Délai jusqu’à 24 h (selon le relecteur)", "livraison ralentie, vulnérabilités exposées"),
    ("Qualité variable", "dépend du relecteur du jour"),
    ("Aucun gate CI/CD", "une PR vulnérable peut être fusionnée"),
    ("Secrets & dépendances", "angles morts de l'œil humain"),
]
yy = 2.6
for t,d in lims:
    rect(s, 7.4, yy+0.05, 0.16, 0.16, fill=RED, shape=MSO_SHAPE.OVAL)
    txt(s, 7.7, yy-0.02, 5.1, 0.7, [[P(t+" — ", 13, NAVY, True, HEAD), P(d, 12, SLATE_T, False, BODY)]])
    yy += 0.85

# ---- 5. OBJECTIFS + CONTRIBUTION ----
s = slide(); header(s, "02 · Objectifs", "La contribution en une phrase", 3)
card = rect(s, 0.8, 2.0, 11.7, 2.1, fill=MIST, line=HAIR, line_w=1.2,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
rect(s, 0.8, 2.0, 0.12, 2.1, fill=GOLD)
txt(s, 1.2, 2.25, 11.0, 1.7, [
    [P("Un ", 21, SLATE_T, False, BODY), P("agent IA autonome", 21, INDIGO, True, HEAD),
     P(" qui revoit chaque Pull Request selon l'", 21, SLATE_T, False, BODY),
     P("OWASP Top 10", 21, NAVY, True, HEAD), P(" en ", 21, SLATE_T, False, BODY),
     P("15 minutes", 21, NAVY, True, HEAD), P(",", 21, SLATE_T, False, BODY)],
    [P("publie les correctifs sur GitHub, et s'exécute ", 21, SLATE_T, False, BODY),
     P("100 % en local", 21, GREEN, True, HEAD),
     P(" — aucun code ne quitte la banque.", 21, SLATE_T, False, BODY)],
])
stats = [("24 h → 15 min", "délai de revue", CYAN_D),
         ("5 / 5", "vulnérabilités détectées", GREEN),
         ("0", "ligne de code ne sort du VPS", INDIGO)]
for i,(big,lab,c) in enumerate(stats):
    x = 0.8 + i*3.95
    rect(s, x, 4.5, 3.7, 1.5, fill=WHITE, line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    txt(s, x, 4.75, 3.7, 0.7, [[P(big, 28, c, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, x, 5.5, 3.7, 0.4, [[P(lab, 12, SLATE, False, BODY)]], align=PP_ALIGN.CENTER)

# ---- 6. POSITIONNEMENT (table) ----
s = slide(); header(s, "02 · Positionnement", "Pourquoi aucune solution existante ne convient", 4)
rows = [
    ("Solution", "Données / inférence", "LLM", "100 % local"),
    ("GitHub CodeQL", "Cloud GitHub", "Non", "Non"),
    ("Snyk", "Code envoyé au cloud", "Limité", "Non"),
    ("SonarQube", "Auto-hébergeable", "Non", "Règles statiques"),
    ("Copilot Autofix", "Cloud", "Oui", "Non"),
    ("BTE Security AI Agent", "Rien ne sort du VPS", "Oui", "OUI"),
]
tb = s.shapes.add_table(len(rows), 4, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.2)).table
tb.columns[0].width=Inches(3.4); tb.columns[1].width=Inches(4.3)
tb.columns[2].width=Inches(1.6); tb.columns[3].width=Inches(2.4)
for r in range(len(rows)):
    for c in range(4):
        cell = tb.cell(r,c); cell.margin_left=Inches(0.12); cell.margin_top=Inches(0.05); cell.margin_bottom=Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]; run = para.add_run(); run.text = rows[r][c]
        run.font.name = HEAD if (r==0 or c==0) else BODY
        run.font.size = Pt(13 if r==0 else 12)
        if r==0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            run.font.color.rgb = WHITE; run.font.bold=True
        elif r==len(rows)-1:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xEA,0xF7,0xF0)
            run.font.color.rgb = NAVY if c<3 else GREEN; run.font.bold = (c==0 or c==3)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r%2 else MIST
            run.font.color.rgb = SLATE_T
            if c==3 and rows[r][c]=="Non":
                run.font.color.rgb = RED
txt(s, 0.8, 6.4, 11.7, 0.4, [[P("Notre différenciateur : ", 13, NAVY, True, HEAD),
    P("100 % local — confidentialité bancaire et conformité BCT.", 13, GREEN, True, HEAD)]])

# ---- 7. ARCHITECTURE ----
s = slide(); header(s, "03 · Architecture", "Vue d'ensemble — 11 conteneurs, 4 couches", 5)
fit_image(s, img("full_architecture.png"), 0.7, 1.8, 8.4, 5.0, caption="Architecture globale du BTE Security AI Agent")
layers = [("Entrée", "nginx — webhook HMAC, reverse proxy", CYAN_D),
          ("IA", "FastAPI + LangGraph + Ollama (7B/14B)", INDIGO),
          ("Données", "PostgreSQL + Redis", NAVY),
          ("Observabilité", "Prometheus · Grafana · AlertManager", GREEN)]
yy = 2.2
for t,d,c in layers:
    rect(s, 9.4, yy, 0.14, 0.9, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 9.7, yy, 3.2, 0.9, [[P(t, 14, NAVY, True, HEAD)],[P(d, 11, SLATE_T, False, BODY)]])
    yy += 1.12

# ---- 8. POURQUOI UN AGENT (6 capabilities) ----
s = slide(); header(s, "04 · IA agentique", "Pourquoi c'est un AGENT, pas un appel LLM", 6)
caps = [
    ("Perception", "événements webhook GitHub", CYAN_D),
    ("Raisonnement", "classe la PR, route le graphe", INDIGO),
    ("Action", "commente, fixe le gate, escalade", NAVY),
    ("État & reprise", "checkpoints LangGraph (PostgreSQL)", GREEN),
    ("Usage d'outils", "19 outils réels (Docker, Prometheus…)", INDIGO),
    ("Auto-exploitation", "scheduler, health digest, gardien disque", CYAN_D),
]
for i,(t,d,c) in enumerate(caps):
    col = i % 3; row = i // 3
    x = 0.8 + col*4.0; y = 2.05 + row*2.25
    rect(s, x, y, 3.7, 2.0, fill=WHITE, line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    dot = rect(s, x+0.3, y+0.32, 0.55, 0.55, fill=c, shape=MSO_SHAPE.OVAL)
    txt(s, x+0.3, y+0.4, 0.55, 0.4, [[P(str(i+1), 18, WHITE, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, x+1.05, y+0.34, 2.5, 0.5, [[P(t, 15, NAVY, True, HEAD)]])
    txt(s, x+0.3, y+1.05, 3.1, 0.8, [[P(d, 12, SLATE_T, False, BODY)]])

# ---- 9. PIPELINE LANGGRAPH ----
s = slide(); header(s, "04 · IA agentique", "Le pipeline de revue — graphe LangGraph", 7)
fit_image(s, img("langgraph_state_graph.png"), 0.7, 1.75, 7.2, 5.05, caption="StateGraph à 9 nœuds avec routage conditionnel")
steps = ["webhook (HMAC + dédup)", "classify  — LLM 7B (~30 s)", "route  — scanners selon la PR",
         "scan  — 5 SAST en parallèle", "analyze — LLM 14B + OWASP", "verdict  — APPROVE / CHANGES / BLOCK",
         "publish — commentaires + gate"]
yy = 2.0
for i,st in enumerate(steps):
    rect(s, 8.2, yy+0.03, 0.28, 0.28, fill=INDIGO if i in (1,4) else CYAN_D, shape=MSO_SHAPE.OVAL)
    txt(s, 8.22, yy+0.03, 0.28, 0.28, [[P(str(i+1), 11, WHITE, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, 8.62, yy, 4.3, 0.5, [[P(st, 12.5, SLATE_T, False, BODY)]])
    yy += 0.66

# ---- 10. DEUX MODÈLES ----
s = slide(); header(s, "04 · IA agentique", "Architecture à deux modèles — décision d'ingénierie", 8)
fit_image(s, img("two_model_architecture.png"), 0.7, 1.8, 7.4, 5.0, caption="7B pour classer · 14B pour analyser")
txt(s, 8.4, 2.05, 4.4, 0.5, [[P("Le bon modèle pour la bonne tâche", 15, NAVY, True, HEAD)]])
pts = [("qwen2.5-coder 7B", "classification rapide ~30 s", INDIGO),
       ("qwen2.5-coder 14B", "analyse de sécurité approfondie", NAVY),
       ("1 seul modèle en RAM", "OLLAMA_MAX_LOADED_MODELS=1", CYAN_D)]
yy = 2.6
for t,d,c in pts:
    rect(s, 8.4, yy+0.04, 0.16, 0.16, fill=c, shape=MSO_SHAPE.OVAL)
    txt(s, 8.7, yy-0.04, 4.2, 0.6, [[P(t, 13, NAVY, True, HEAD)],[P(d, 11.5, SLATE_T, False, BODY)]])
    yy += 0.85
rect(s, 8.4, 5.55, 4.4, 1.0, fill=RGBColor(0xEA,0xF7,0xF0), line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 8.6, 5.7, 4.0, 0.8, [[P("−50 % ", 22, GREEN, True, HEAD),
    P("de temps total", 14, NAVY, True, HEAD)],
    [P("fusion de 2 appels 14B en un seul", 11.5, SLATE_T, False, BODY)]])

# ---- 11. ANTI-HALLUCINATION ----
s = slide(); header(s, "04 · IA agentique", "Confiance : le système anti-hallucination", 9)
fit_image(s, img("anti_hallucination_layers.png"), 0.7, 1.8, 6.4, 5.0, caption="6 couches de protection empilées")
fit_image(s, img("annotated_diff_format.png"), 7.2, 1.9, 5.7, 3.1, caption="Diff annoté — parser anti-hallucination")
txt(s, 7.2, 5.45, 5.7, 1.2, [
    [P("Le doute n°1 du jury sur les LLM : ", 13, NAVY, True, HEAD),
     P("« peut-on lui faire confiance ? »", 13, SLATE_T, False, BODY, True)],
    [P("temperature=0 · num_ctx maîtrisé · garde sans-outil · prompt anti-hallucination · "
       "observation verbatim · parser qui élimine les lignes inexistantes.", 12, SLATE_T, False, BODY)],
])

# ---- 12. CHAT OPÉRATIONNEL ----
s = slide(); header(s, "04 · IA agentique", "Seconde surface agentique : le chat opérationnel", 10)
fit_image(s, img("react_loop_architecture.png"), 0.7, 1.8, 7.3, 5.0, caption="Boucle ReAct personnalisée — 19 outils de monitoring")
txt(s, 8.3, 2.1, 4.6, 2.0, [
    [P("« Quel est l'usage CPU ? »", 15, INDIGO, True, HEAD)],
    [P("L'agent appelle un vrai outil ", 13, SLATE_T, False, BODY),
     P("(vps_status)", 12, NAVY, True, MONO),
     P(" et répond avec la valeur réelle — jamais inventée.", 13, SLATE_T, False, BODY)],
])
for i,(cat,n) in enumerate([("VPS / Host","5"),("Docker","6"),("Prometheus","3"),("Ollama / Redis","2"),("Artifacts / BD","3")]):
    y = 3.6 + i*0.55
    rect(s, 8.3, y, 0.5, 0.42, fill=CYAN_D, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 8.3, y+0.03, 0.5, 0.36, [[P(n, 14, WHITE, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, 8.95, y+0.06, 3.8, 0.4, [[P(cat, 13, SLATE_T, False, BODY)]])

# ---- 13. PIPELINE SÉCURITÉ ----
s = slide(); header(s, "05 · Sécurité", "Défense en profondeur : 5 scanners + LLM", 11)
fit_image(s, img("sast_tools_overview.png"), 0.7, 1.8, 7.6, 5.0, caption="5 outils SAST en parallèle — couverture OWASP Top 10 2025")
tools = [("Trivy","CVE & images"),("Gitleaks","secrets"),("Semgrep","OWASP code"),
         ("Checkov","IaC"),("OSV-Scanner","dépendances")]
yy = 2.1
for t,d in tools:
    chip(s, 8.6, yy, 1.8, t, CYAN_D)
    txt(s, 10.55, yy+0.02, 2.3, 0.4, [[P(d, 12, SLATE_T, False, BODY)]])
    yy += 0.55
rect(s, 8.6, 5.05, 4.25, 1.5, fill=RGBColor(0xF3,0xF0,0xFD), line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 8.8, 5.2, 3.9, 1.3, [
    [P("Complémentarité", 14, INDIGO, True, HEAD)],
    [P("Le LLM a capté des secrets en dur que Gitleaks avait manqués — "
       "les scanners et l'IA se renforcent.", 12, SLATE_T, False, BODY)],
])

# ---- 14. BÉNÉFICES (table AS-IS / TO-BE) ----
s = slide(); header(s, "05 · Sécurité", "Bénéfices : AS-IS vs TO-BE", 12)
rows = [
    ("Critère", "AS-IS — manuel", "TO-BE — agent IA"),
    ("Délai de revue", "≤ 24 h, selon le relecteur", "15–25 min"),
    ("Couverture", "variable selon le relecteur", "OWASP Top 10 systématique"),
    ("Gate CI/CD", "aucun", "APPROVE / REQUEST_CHANGES / BLOCK"),
    ("Confidentialité", "—", "100 % local (conformité BCT)"),
    ("Traçabilité", "historique GitHub", "base PostgreSQL + dashboard RSSI"),
]
tb = s.shapes.add_table(len(rows), 3, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.4)).table
tb.columns[0].width=Inches(2.9); tb.columns[1].width=Inches(4.0); tb.columns[2].width=Inches(4.8)
for r in range(len(rows)):
    for c in range(3):
        cell = tb.cell(r,c); cell.margin_left=Inches(0.14); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_top=Inches(0.06); cell.margin_bottom=Inches(0.06)
        para = cell.text_frame.paragraphs[0]; run = para.add_run(); run.text = rows[r][c]
        run.font.name = HEAD if (r==0 or c==0) else BODY
        run.font.size = Pt(13 if r==0 else 12.5)
        if r==0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY; run.font.color.rgb = WHITE; run.font.bold=True
            if c==2: cell.fill.fore_color.rgb = GREEN
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r%2 else MIST
            run.font.color.rgb = SLATE_T
            if c==0: run.font.color.rgb = NAVY; run.font.bold=True
            if c==2: run.font.color.rgb = RGBColor(0x12,0x7A,0x46); run.font.bold=True

# ---- 15. DÉMO ----
s = slide(); header(s, "06 · Démonstration", "PR #18 — revue en conditions réelles", 13)
fit_image(s, img("pr_comment-security-review_0.png"), 0.7, 1.85, 7.0, 4.95, card=True,
          caption="Revue de sécurité publiée automatiquement sur la Pull Request")
txt(s, 8.0, 2.1, 4.9, 0.5, [[P("Ce que fait l'agent, sans intervention", 15, NAVY, True, HEAD)]])
flow = [("Reçoit le webhook", CYAN_D), ("Lance 5 scanners + LLM", INDIGO),
        ("Poste les commentaires inline", NAVY), ("Suggère les correctifs", CYAN_D)]
yy=2.7
for t,c in flow:
    rect(s, 8.0, yy+0.04, 0.16,0.16, fill=c, shape=MSO_SHAPE.OVAL)
    txt(s, 8.3, yy-0.02, 4.6, 0.4, [[P(t, 13, SLATE_T, False, BODY)]])
    yy += 0.6
chip(s, 8.0, 5.35, 4.9, "VERDICT :  BLOCK", RED, h=0.6)
txt(s, 8.0, 6.05, 4.9, 0.5, [[P("→ gate commit rouge, fusion bloquée", 12, SLATE, False, BODY, True)]])

# ---- 16. RÉSULTATS ----
s = slide(); header(s, "06 · Résultats", "Preuve : détection mesurée", 14)
rect(s, 0.8, 2.1, 5.6, 3.0, fill=RGBColor(0xEA,0xF7,0xF0), line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
txt(s, 0.8, 2.5, 5.6, 1.2, [[P("5 / 5", 64, GREEN, True, HEAD)]], align=PP_ALIGN.CENTER)
txt(s, 0.8, 3.9, 5.6, 0.9, [[P("vulnérabilités plantées détectées", 16, NAVY, True, HEAD)],
    [P("sur la PR #18 de démonstration", 13, SLATE, False, BODY)]], align=PP_ALIGN.CENTER)
vulns = ["Injection SQL", "Secrets en dur", "Path traversal", "Injection de commande", "Hachage faible (MD5)"]
txt(s, 6.8, 2.15, 6.0, 0.5, [[P("Les 5 failles couvertes", 15, NAVY, True, HEAD)]])
yy=2.75
for v in vulns:
    rect(s, 6.8, yy+0.04, 0.18,0.18, fill=GREEN, shape=MSO_SHAPE.OVAL)
    txt(s, 7.1, yy-0.02, 5.6, 0.4, [[P(v, 14, SLATE_T, False, BODY)]])
    yy += 0.55
txt(s, 6.8, 5.7, 6.0, 0.8, [[P("+ corpus de 8 revues réelles", 14, NAVY, True, HEAD),
    P("  persistées dans PostgreSQL.", 13, SLATE_T, False, BODY)],
    [P("Périmètre assumé : jeu contrôlé, pas un benchmark à l'échelle.", 12, SLATE, False, BODY, True)]])

# ---- 17. ROBUSTESSE ----
s = slide(); header(s, "07 · Robustesse", "Maturité production : observabilité & incidents", 15)
fit_image(s, img("grafana_pr_reviews_dashboard.png"), 0.7, 1.85, 7.2, 4.4, caption="Dashboard Grafana — PR Security Reviews")
txt(s, 8.2, 2.05, 4.7, 0.5, [[P("3 incidents réels survécus", 15, NAVY, True, HEAD)]])
inc = [("Urgence disque", "gardien auto-réparateur à 90 %"),
       ("AlertManager cassé", "détecté par l'observabilité"),
       ("Bug template Phi-4", "corrigé et documenté")]
yy=2.65
for t,d in inc:
    rect(s, 8.2, yy+0.04, 0.16,0.16, fill=AMBER, shape=MSO_SHAPE.OVAL)
    txt(s, 8.5, yy-0.04, 4.4, 0.6, [[P(t, 13, NAVY, True, HEAD)],[P(d, 11.5, SLATE_T, False, BODY)]])
    yy += 0.82
txt(s, 8.2, 5.3, 4.7, 1.0, [[P("Prometheus · VictoriaMetrics · Grafana · AlertManager", 12, SLATE_T, False, BODY)],
    [P("28 métriques · 15 règles d'alerte · rétention 90 j", 12, CYAN_D, True, HEAD)]])

# ---- 18. LIMITES & PERSPECTIVES ----
s = slide(); header(s, "07 · Perspectives", "Limites assumées & perspectives", 16)
rect(s, 0.8, 2.0, 5.7, 4.4, fill=MIST, line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 1.05, 2.2, 5.2, 0.5, [[P("Limites", 16, RED, True, HEAD)]])
for i,t in enumerate(["Inférence CPU : ~3 tokens/s (15–25 min/PR)",
                      "Durcissement : ports ouverts, mot de passe Redis",
                      "Corpus d'évaluation restreint (jeu contrôlé)"]):
    txt(s, 1.05, 2.8+i*0.95, 5.2, 0.8, [[P("•  ", 13, RED, True, HEAD), P(t, 13, SLATE_T, False, BODY)]])
rect(s, 6.9, 2.0, 5.6, 4.4, fill=RGBColor(0xEC,0xFB,0xFD), line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 7.15, 2.2, 5.1, 0.5, [[P("Perspectives", 16, GREEN, True, HEAD)]])
for i,t in enumerate(["Migration GPU : 30–50 tok/s → 2–5 min/PR",
                      "Passage à l'échelle de tous les dépôts BTE",
                      "Plan de durcissement production"]):
    txt(s, 7.15, 2.8+i*0.95, 5.1, 0.8, [[P("•  ", 13, GREEN, True, HEAD), P(t, 13, SLATE_T, False, BODY)]])

# ---- 19. CONCLUSION ----
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=CREAM)
rect(s, 0.9, 1.35, 11.5, 0.02, fill=NAVY); rect(s, 0.9, 1.35, 2.4, 0.02, fill=GOLD)
txt(s, 1.0, 1.55, 11.3, 0.5, [[P("CONCLUSION", 14, CYAN_D, True, BODY)]])
txt(s, 1.0, 2.2, 11.3, 2.2, [
    [P("Un agent IA agentique, déployé en production à la BTE,", 26, NAVY, True, HEAD)],
    [P("qui sécurise chaque Pull Request ", 26, NAVY, True, HEAD),
     P("sans qu'aucun code ne quitte la banque.", 26, CYAN_D, True, HEAD)],
])
for i,(big,lab) in enumerate([("24 h → 15 min","délai"),("5/5","détection"),("100 % local","conformité BCT")]):
    x = 1.0 + i*3.9
    rect(s, x, 4.45, 3.6, 1.3, fill=WHITE, line=HAIR, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, 4.45, 0.10, 1.3, fill=GOLD)
    txt(s, x, 4.7, 3.6, 0.6, [[P(big, 22, CYAN_D, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, x, 5.35, 3.6, 0.4, [[P(lab, 12, SLATE, False, BODY)]], align=PP_ALIGN.CENTER)
txt(s, 1.0, 6.2, 11.3, 0.6, [[P("Plateforme remise à la BTE  ·  Merci de votre attention", 16, GOLD, True, HEAD)]])
rect(s, 0.9, 6.95, 11.5, 0.02, fill=HAIR)

prs.save(OUT)
print("OK ->", OUT, "| slides:", len(prs.slides._sldIdLst))
